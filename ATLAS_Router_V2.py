import math
import json
import uuid
import time
import hashlib
import statistics
from copy import deepcopy
from dataclasses import dataclass, field, asdict
from typing import Dict, List, Any, Tuple, Optional, Callable
from datetime import datetime, timezone
def _imul(a: int, b: int) -> int:
    a &= 0xFFFFFFFF
    b &= 0xFFFFFFFF
    return (a * b) & 0xFFFFFFFF

class RNG:
    """Seeded deterministic PRNG with normal / gamma / beta sampling."""
    def __init__(self, seed: int = 42):
        self.state = seed & 0xFFFFFFFF

    def next_float(self) -> float:
        self.state = (self.state + 0x6D2B79F5) & 0xFFFFFFFF
        t = self.state
        t = _imul(t ^ (t >> 15), t | 1)
        t ^= (t + _imul(t ^ (t >> 7), t | 61)) & 0xFFFFFFFF
        t &= 0xFFFFFFFF
        return ((t ^ (t >> 14)) & 0xFFFFFFFF) / 4294967296.0

    def _normal(self) -> float:
        import math
        u1 = max(self.next_float(), 1e-12)
        u2 = self.next_float()
        return math.sqrt(-2.0 * math.log(u1)) * math.cos(2.0 * math.pi * u2)

    def gamma(self, k: float) -> float:
        import math
        if k <= 0:
            raise ValueError("Gamma shape must be > 0.")
        if k < 1.0:
            u = max(self.next_float(), 1e-12)
            return self.gamma(k + 1.0) * (u ** (1.0 / k))
        d = k - 1.0 / 3.0
        c = 1.0 / math.sqrt(9.0 * d)
        while True:
            x = self._normal()
            v = (1.0 + c * x) ** 3
            if v <= 0:
                continue
            u = self.next_float()
            if u < 1.0 - 0.0331 * (x ** 4):
                return d * v
            if math.log(max(u, 1e-12)) < 0.5 * x * x + d * (1.0 - v + math.log(v)):
                return d * v

    def beta(self, alpha: float, beta_param: float) -> float:
        if alpha <= 0 or beta_param <= 0:
            raise ValueError("Beta parameters must be > 0.")
        x = self.gamma(alpha)
        y = self.gamma(beta_param)
        if x + y == 0:
            return 0.5
        return x / (x + y)

GLOBAL_RNG = RNG(42)



# Each is wrapped so the router degrades gracefully when a library is absent.
import os
import mimetypes

def _try_import(name):
    try:
        return __import__(name)
    except ImportError:
        return None

fitz = _try_import("fitz")          # PyMuPDF  -> PDFs
PIL = _try_import("PIL")            # Pillow   -> images
openpyxl = _try_import("openpyxl")  # Excel
pptx_mod = _try_import("pptx")      # python-pptx -> PowerPoint
mutagen = _try_import("mutagen")    # audio metadata

try:
    from PIL import Image as _PILImage
except ImportError:
    _PILImage = None

import subprocess
import shutil

def _ffprobe_available() -> bool:
    return shutil.which("ffprobe") is not None

ARTIFACT_LIBS = {
    "pymupdf": fitz is not None,
    "pillow": _PILImage is not None,
    "openpyxl": openpyxl is not None,
    "python-pptx": pptx_mod is not None,
    "mutagen": mutagen is not None,
    "ffprobe": _ffprobe_available(),
}

ROUTER_VERSION = "atlas-router-v2.0.0"
PARSER_VERSION = "parser-v2.0.0"
POLICY_VERSION = "policy-v2.0.0"
SCORING_VERSION = "scoring-v2.0.0"
CALIBRATION_VERSION = "calibration-v2.0.0"
REGISTRY_VERSION = "registry-v2.0.0-2026-06-snapshot"
TELEMETRY_SNAPSHOT_VERSION = "telemetry-snap-2026-06-27"
SUPPORTED_FORMATS = [
    "text", "image", "pdf", "audio", "video", "spreadsheet", "presentation"
]

CONFIDENCE_ABSTAIN_THRESHOLD = 0.20
CONFIDENCE_ESCALATE_THRESHOLD = 0.20
CONFIDENCE_HIGH_THRESHOLD = 0.45

Models = {
    'OpenAI': [
        'GPT-5.5', 'GPT-5.5-Thinking', 'GPT-5.5-Instant', 'GPT-5.4', 
        'GPT-5.4-Mini', 'GPT-5', 'GPT-4o', 'GPT-4o-Mini', 'o4-Mini', 'o3'
    ],
    'Anthropic': [
        'Claude-Opus-4.8', 'Claude-Opus-4.7', 'Claude-Opus-4.6', 
        'Claude-Opus-4.5', 'Claude-Sonnet-4.6', 'Claude-Sonnet-4.5', 
        'Claude-Sonnet-4', 'Claude-Haiku-4.5', 'Claude-Haiku-4', 
        'Claude-3.5-Sonnet'
    ],
    'Google': [
        'Gemini-3.1-Pro', 'Gemini-3-Pro', 'Gemini-3.5-Flash', 'Gemini-2.5-Pro', 
        'Gemini-2.5-Flash', 'Gemini-2.0-Pro', 'Gemini-2.0-Flash', 'Gemma-4', 'Gemma-3'
    ],
    'Meta': [
        'Llama-4-Maverick', 'Llama-4-Scout', 'Llama-3.3-70B', 
        'Llama-3.1-405B', 'Llama-3.1-70B', 'Llama-3.2-90B-Vision', 
        'Llama-3.2-11B-Vision', 'Llama-3.2-3B', 'Llama-3.2-1B'
    ],
    'DeepSeek': [
        'DeepSeek-V4-Pro', 'DeepSeek-R1', 'DeepSeek-R1-Distill', 
        'DeepSeek-V3.2', 'DeepSeek-V3', 'DeepSeek-V2.5'
    ],
    'Alibaba (Qwen)': [
        'Qwen-3.7-Plus', 'Qwen-3.5', 'Qwen-3', 'Qwen-3-Coder', 
        'Qwen-2.5-Coder', 'Qwen-2.5-72B', 'Qwen-2.5-32B', 'Qwen-2.5-7B'
    ],
    'Mistral AI': [
        'Mistral-Large-3', 'Mistral-Large-2', 'Mistral-Medium-3.5', 
        'Mistral-Small', 'Mistral-7B', 'Mixtral-8x22B', 'Mixtral-8x7B', 'Codestral'
    ],
    'xAI': [
        'Grok-4.3', 'Grok-4', 'Grok-3', 'Grok-2'
    ],
    'Microsoft': [
        'Phi-4', 'Phi-4-Multimodal', 'Phi-4-Mini', 'Phi-3-Medium', 
        'Phi-3-Small', 'Phi-3-Mini'
    ],
    'Cohere': [
        'Command-A', 'Command-R-Plus', 'Command-R'
    ],
    'Moonshot AI': [
        'Kimi-K2.7-Code', 'Kimi-K2.5', 'Kimi-K2', 'Kimi-K1'
    ],
    'Z.ai / GLM': [
        'GLM-5.2', 'GLM-5', 'GLM-4-Plus', 'GLM-4-Air', 'GLM-4'
    ],
    'NVIDIA': [
        'Nemotron-Ultra', 'Nemotron-70B', 'Nemotron-4'
    ],
    'Amazon': [
        'Nova-Premier', 'Nova-Pro', 'Nova-Lite', 'Nova-Micro'
    ],
    'IBM': [
        'Granite-3.5', 'Granite-3'
    ],
    'AI21': [
        'Jamba-Large', 'Jamba-Mini'
    ],
    '01.AI': [
        'Yi-Large', 'Yi-Lightning', 'Yi-34B'
    ],
    'MiniMax': [
        'MiniMax-M3', 'MiniMax-M2', 'MiniMax-M1'
    ],
    'Baidu': [
        'ERNIE-4.5', 'ERNIE-X1'
    ]
}

from model_registry import MODEL_REGISTRY


@dataclass
class ArtifactProfile:
    format: str
    page_count: Optional[int] = None
    text_density: Optional[float] = None
    scan_likelihood: Optional[float] = None
    handwriting_likelihood: Optional[float] = None
    table_density: Optional[float] = None
    chart_density: Optional[float] = None
    detected_language: Optional[str] = None
    audio_duration_sec: Optional[int] = None
    audio_quality: Optional[float] = None
    video_duration_sec: Optional[int] = None
    spreadsheet_complexity: Optional[float] = None
    presentation_complexity: Optional[float] = None
    # ---- New fields for the file-driven pipeline ----
    source_path: Optional[str] = None
    file_size_bytes: Optional[int] = None
    extracted_text_preview: Optional[str] = None   
    inferred_topic: Optional[str] = None       
    extraction_method: str = "heuristic"         



@dataclass
class StructuredSemanticParse:
    primary_family: str
    secondary_families: List[str]
    required_stages: List[str]
    workflow_graph: List[Dict[str, Any]]
    domain: str
    risk_tier: str
    risk_type: str
    expected_output: str
    ambiguity_score: float
    actionability: str
    document_type: str
    decomposition_needed: bool
    needs_verification: bool
    parser_confidence: float = 0.75
    reason_summary: str = ""


@dataclass
class RequestConstraints:
    max_latency_ms: Optional[int] = None
    max_cost_usd: Optional[float] = None
    allowed_providers: List[str] = field(default_factory=list)
    disallowed_providers: List[str] = field(default_factory=list)
    allowed_tiers: List[str] = field(default_factory=list)
    no_open_weight: bool = False
    no_web_access: bool = False
    min_confidence: float = 0.20
    required_region: Optional[str] = None
    privacy_class: str = "standard"
    must_use_single_model: bool = False
    mandatory_verifier: bool = False


@dataclass
class TenantContext:
    tenant_id: str = "default"
    tenant_name: str = "default"
    policy_overlay: Dict[str, Any] = field(default_factory=dict)
    budget_remaining_usd: Optional[float] = None
    preferred_profiles: List[str] = field(default_factory=list)
    allowed_models: List[str] = field(default_factory=list)


@dataclass
class TaskFeatures:
    raw_prompt: str
    input_formats: List[str]
    estimated_tokens: int
    estimated_output_tokens: int
    artifacts: List[ArtifactProfile]
    primary_family: str
    secondary_families: List[str]
    required_stages: List[str]
    workflow_graph: List[Dict[str, Any]]
    complexity: str
    domain: str
    risk_tier: str
    risk_type: str
    required_formats: List[str]
    requires_json: bool
    requires_function_calling: bool
    requires_web_search: bool
    requires_ocr: bool
    requires_citations: bool
    requires_verifier: bool
    min_context_window: int
    expected_output: str
    ambiguity_score: float
    safety_sensitive: bool
    actionability: str
    document_type: str
    decomposition_needed: bool
    workflow_profile: str
    parser_confidence: float
    conflict_flags: List[str] = field(default_factory=list)
    extracted_text_summary: Optional[str] = None
    detected_languages: List[str] = field(default_factory=list)
    total_file_size_bytes: int = 0
    request_constraints: RequestConstraints = field(default_factory=RequestConstraints)
    tenant_context: TenantContext = field(default_factory=TenantContext)
    


@dataclass
class PolicyDecision:
    allowed: bool
    must_escalate: bool
    must_abstain: bool
    restricted_to_tiers: List[str] = field(default_factory=list)
    restricted_to_models: List[str] = field(default_factory=list)
    restricted_to_providers: List[str] = field(default_factory=list)
    require_verifier_types: List[str] = field(default_factory=list)
    deny_reason: Optional[str] = None
    notes: List[str] = field(default_factory=list)


@dataclass
class StageRoute:
    stage_id: int
    stage_name: str
    selected_model: Optional[str]
    fallback_models: List[str]
    verifier_models: List[str]
    stage_confidence: float
    expected_latency_ms: float
    expected_cost_usd: float
    explanation: str


@dataclass
class ExecutionPlan:
    plan_id: str
    plan_type: str  # "single_model" or "multi_stage"
    selected_model: Optional[str]
    stage_routes: List[StageRoute]
    fallback_models: List[str]
    verifier_models: List[str]
    expected_latency_ms: float
    expected_cost_usd: float
    expected_quality: float
    confidence: float
    utility: float
    confidence_margin: float
    profile_used: str
    explanation: Dict[str, Any]
    trace: Dict[str, Any]


@dataclass
class RoutingDecision:
    selected_plan: Optional[ExecutionPlan]
    abstain: bool
    escalate_to_human: bool
    decision_record: Dict[str, Any]


TASK_FAMILY_SIGNAL_MAP: Dict[str, Dict[str, float]] = {
    "coding": {
        "perf:coding": 0.50, "perf:agentic_tasks": 0.15,
        "bench:swe_bench": 0.15, "bench:humaneval": 0.10,
        "perf:instruction_following": 0.10
    },
    "reasoning": {
        "perf:reasoning": 0.42, "perf:scientific_reasoning": 0.18,
        "bench:gpqa": 0.20, "bench:mmlu": 0.10,
        "perf:instruction_following": 0.10
    },
    "mathematics": {
        "perf:mathematics": 0.50, "bench:aime": 0.25,
        "perf:reasoning": 0.15, "bench:gpqa": 0.10
    },
    "chat": {
        "perf:creative_writing": 0.20, "perf:instruction_following": 0.40,
        "domain:general": 0.20, "domain:customer_support": 0.20
    },
    "vision": {
        "perf:vision_understanding": 0.50, "bench:mmmu": 0.25,
        "perf:ocr": 0.10, "perf:table_understanding": 0.15
    },
    "ocr": {
        "perf:ocr": 0.55, "bench:docvqa": 0.20,
        "perf:table_understanding": 0.15, "perf:document_qa": 0.10
    },
    "document_qa": {
        "perf:document_qa": 0.40, "perf:long_context": 0.20,
        "bench:docvqa": 0.15, "perf:summarization": 0.10,
        "perf:table_understanding": 0.15
    },
    "summarization": {
        "perf:summarization": 0.50, "perf:long_context": 0.20,
        "perf:instruction_following": 0.20, "perf:document_qa": 0.10
    },
    "translation": {
        "perf:translation": 0.70, "perf:instruction_following": 0.30
    },
    "agent": {
        "perf:agentic_tasks": 0.45, "perf:instruction_following": 0.20,
        "perf:reasoning": 0.20, "perf:coding": 0.15
    },
    "audio": {
        "perf:audio_understanding": 0.60, "perf:summarization": 0.20,
        "perf:instruction_following": 0.20
    },
    "creative": {
        "perf:creative_writing": 0.40, "perf:instruction_following": 0.25,
        "perf:reasoning": 0.15, "perf:vision_understanding": 0.20
    },
    "research": {
        "perf:reasoning": 0.30, "perf:long_context": 0.20,
        "perf:instruction_following": 0.20, "perf:summarization": 0.15,
        "perf:scientific_reasoning": 0.15
    },
    "education": {
        "perf:instruction_following": 0.35, "perf:reasoning": 0.25,
        "perf:creative_writing": 0.20, "perf:scientific_reasoning": 0.20
    },
    "data_analysis": {
        "perf:reasoning": 0.30, "perf:coding": 0.25,
        "perf:table_understanding": 0.20, "perf:mathematics": 0.15,
        "perf:instruction_following": 0.10
    }
}

DOMAIN_SAFETY_MAP: Dict[str, str] = {
    "medical": "medical", "legal": "legal", "finance": "finance",
    "security": "cybersecurity", "sensitive": "sensitive", "general": "general",
    "software": "general", "science": "general", "research": "general",
    "education": "general", "mathematics": "general",
    "customer_support": "general"
}

DOMAIN_EXPERTISE_MAP: Dict[str, str] = {
    "medical": "medical", "legal": "legal", "finance": "finance",
    "security": "software", "sensitive": "general", "software": "software",
    "science": "science", "research": "research", "education": "education",
    "general": "general", "mathematics": "mathematics",
    "customer_support": "customer_support"
}


EXTENSION_FORMAT_MAP: Dict[str, str] = {
    ".pdf": "pdf",
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".gif": "image", ".bmp": "image", ".tiff": "image", ".webp": "image",
    ".mp3": "audio", ".wav": "audio", ".m4a": "audio",
    ".flac": "audio", ".ogg": "audio", ".aac": "audio",
    ".mp4": "video", ".mov": "video", ".mkv": "video",
    ".avi": "video", ".webm": "video",
    ".xlsx": "spreadsheet", ".xls": "spreadsheet", ".csv": "spreadsheet",
    ".pptx": "presentation", ".ppt": "presentation",
    ".txt": "text", ".md": "text",
}

def detect_format(path: str) -> Optional[str]:
    """Determine the canonical SUPPORTED_FORMATS label for a file path.

    Falls back to MIME-type sniffing when the extension is unknown.
    Returns None if the format cannot be mapped (caller decides how to handle).
    """
    ext = os.path.splitext(path)[1].lower()
    if ext in EXTENSION_FORMAT_MAP:
        return EXTENSION_FORMAT_MAP[ext]
    mime, _ = mimetypes.guess_type(path)
    if mime:
        if mime.startswith("image/"):
            return "image"
        if mime.startswith("audio/"):
            return "audio"
        if mime.startswith("video/"):
            return "video"
        if mime == "application/pdf":
            return "pdf"
        if "spreadsheet" in mime or mime == "text/csv":
            return "spreadsheet"
        if "presentation" in mime:
            return "presentation"
        if mime.startswith("text/"):
            return "text"
    return None


def _inspect_pdf(path: str, profile: "ArtifactProfile") -> None:
    """Read PDF structure with PyMuPDF; derive scan/text/table signals."""
    if fitz is None:
        return
    try:
        doc = fitz.open(path)
        profile.page_count = doc.page_count
        # Sample the first page to estimate text density vs. scanned image.
        first = doc.load_page(0)
        text = first.get_text("text") or ""
        images = first.get_images(full=True)
        char_count = len(text.strip())
        # Heuristic: lots of extractable text -> native; little text + images -> scan.
        profile.text_density = min(1.0, char_count / 1500.0)
        profile.scan_likelihood = 0.85 if (char_count < 100 and images) else 0.10
        profile.handwriting_likelihood = 0.05
        # Crude table signal: many tab/newline-aligned columns.
        profile.table_density = 0.70 if text.count("\t") > 10 else 0.10
        profile.detected_language = "en"  # replace with langdetect in prod
        doc.close()
    except Exception:
        pass

def _inspect_image(path: str, profile: "ArtifactProfile") -> None:
    if _PILImage is None:
        return
    try:
        img = _PILImage.open(path)
        w, h = img.size
        # Screenshots tend to be wide; receipts/scans tend to be tall.
        profile.scan_likelihood = 0.60 if h > w else 0.20
        profile.handwriting_likelihood = 0.10
        profile.table_density = 0.10
        profile.chart_density = 0.10
        profile.detected_language = "en"
    except Exception:
        pass

def _inspect_media_ffprobe(path: str, profile: "ArtifactProfile", is_video: bool) -> bool:
    """Use ffprobe to read duration and stream info. Returns True on success."""
    if not _ffprobe_available():
        return False
    try:
        out = subprocess.run(
            ["ffprobe", "-v", "quiet", "-print_format", "json",
             "-show_format", "-show_streams", path],
            capture_output=True, text=True, timeout=20
        )
        meta = json.loads(out.stdout or "{}")
        duration = float(meta.get("format", {}).get("duration", 0)) or None
        has_audio = any(s.get("codec_type") == "audio" for s in meta.get("streams", []))
        if is_video:
            profile.video_duration_sec = int(duration) if duration else None
            profile.audio_quality = 0.70 if has_audio else 0.0
            profile.chart_density = 0.20
        else:
            profile.audio_duration_sec = int(duration) if duration else None
            profile.audio_quality = 0.75
        profile.detected_language = "en"
        return True
    except Exception:
        return False

def _inspect_audio(path: str, profile: "ArtifactProfile") -> None:
    if _inspect_media_ffprobe(path, profile, is_video=False):
        return
    if mutagen is not None:
        try:
            audio = mutagen.File(path)
            if audio is not None and audio.info is not None:
                profile.audio_duration_sec = int(audio.info.length)
                profile.audio_quality = 0.75
                profile.detected_language = "en"
        except Exception:
            pass

def _inspect_video(path: str, profile: "ArtifactProfile") -> None:
    _inspect_media_ffprobe(path, profile, is_video=True)

def _inspect_spreadsheet(path: str, profile: "ArtifactProfile") -> None:
    if openpyxl is None:
        return
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        n_sheets = len(wb.sheetnames)
        # Complexity rises with multiple sheets / large used ranges.
        max_cells = 0
        for ws in wb.worksheets:
            max_cells = max(max_cells, (ws.max_row or 0) * (ws.max_column or 0))
        profile.spreadsheet_complexity = min(1.0, 0.3 + 0.2 * n_sheets + max_cells / 50000.0)
        profile.table_density = 0.90
        profile.detected_language = "en"
        wb.close()
    except Exception:
        pass

def _inspect_presentation(path: str, profile: "ArtifactProfile") -> None:
    if pptx_mod is None:
        return
    try:
        from pptx import Presentation
        prs = Presentation(path)
        n_slides = len(prs.slides)
        profile.presentation_complexity = min(1.0, 0.3 + n_slides / 40.0)
        profile.chart_density = 0.55
        profile.detected_language = "en"
    except Exception:
        pass


def inspect_artifacts(
    input_formats: Optional[List[str]] = None,
    prompt: str = "",
    artifact_hints: Optional[List[Dict[str, Any]]] = None,
    files: Optional[List[str]] = None,
) -> List[ArtifactProfile]:
    """Produce one ArtifactProfile per input.

    Two modes:
      1. FILE MODE (preferred): pass `files=[...]`. Real metadata is read from
         each file using PyMuPDF/Pillow/FFprobe/openpyxl/python-pptx when those
         libraries are available, with graceful degradation otherwise.
      2. LEGACY MODE: pass `input_formats=[...]`. Falls back to prompt-keyword
         heuristics. Preserved for backward compatibility and synthetic tests.

    Caller-supplied `artifact_hints` always win over both file reads and
    heuristics (ground truth injection from an upstream system).
    """
    artifact_hints = artifact_hints or []
    profiles: List[ArtifactProfile] = []

    # ---------------- FILE MODE ----------------
    if files:
        hints_by_format = {h.get("format"): h for h in artifact_hints if "format" in h}
        for path in files:
            fmt = detect_format(path)
            if fmt is None:
                # Unmappable file: record as text with a note rather than crash.
                profiles.append(ArtifactProfile(format="text"))
                continue
            profile = ArtifactProfile(format=fmt)
            try:
                if fmt == "pdf":
                    _inspect_pdf(path, profile)
                elif fmt == "image":
                    _inspect_image(path, profile)
                elif fmt == "audio":
                    _inspect_audio(path, profile)
                elif fmt == "video":
                    _inspect_video(path, profile)
                elif fmt == "spreadsheet":
                    _inspect_spreadsheet(path, profile)
                elif fmt == "presentation":
                    _inspect_presentation(path, profile)
            except Exception:
                pass
            # Lightweight semantic extraction.
            try:
                profile.source_path = path
                profile.file_size_bytes = os.path.getsize(path)
                profile.extraction_method = "library"
                preview = extract_text_preview(path, fmt)
                profile.extracted_text_preview = preview
                profile.inferred_topic = infer_topic(preview)
            except Exception:
                pass
            # Apply hint overrides (ground truth wins).

            h = hints_by_format.get(fmt, {})
            for k, v in h.items():
                if k != "format" and hasattr(profile, k):
                    setattr(profile, k, v)
            profiles.append(profile)
        return profiles

    # ---------------- LEGACY HEURISTIC MODE ----------------
    input_formats = input_formats or ["text"]
    p = prompt.lower()
    hints_by_format = {h.get("format"): h for h in artifact_hints if "format" in h}
    for fmt in sorted(set(input_formats)):
        h = hints_by_format.get(fmt, {})
        profile = ArtifactProfile(format=fmt)
        if fmt == "pdf":
            profile.page_count = h.get("page_count", 30 if "contract" in p else 10)
            profile.text_density = h.get("text_density",
                0.25 if any(k in p for k in ["scanned", "scan", "invoice"]) else 0.80)
            profile.scan_likelihood = h.get("scan_likelihood",
                0.85 if any(k in p for k in ["scanned", "scan", "invoice"]) else 0.10)
            profile.handwriting_likelihood = h.get("handwriting_likelihood",
                0.50 if "handwritten" in p else 0.05)
            profile.table_density = h.get("table_density",
                0.70 if any(k in p for k in ["invoice", "statement", "table"]) else 0.10)
            profile.chart_density = h.get("chart_density",
                0.55 if any(k in p for k in ["chart", "graph"]) else 0.10)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "image":
            profile.scan_likelihood = h.get("scan_likelihood",
                0.80 if any(k in p for k in ["extract text", "scan", "ocr"]) else 0.10)
            profile.handwriting_likelihood = h.get("handwriting_likelihood",
                0.60 if "handwritten" in p else 0.10)
            profile.table_density = h.get("table_density",
                0.65 if any(k in p for k in ["invoice", "receipt", "table"]) else 0.10)
            profile.chart_density = h.get("chart_density",
                0.50 if "chart" in p else 0.05)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "audio":
            profile.audio_duration_sec = h.get("audio_duration_sec",
                900 if any(k in p for k in ["call", "meeting"]) else 120)
            profile.audio_quality = h.get("audio_quality", 0.75)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "video":
            profile.video_duration_sec = h.get("video_duration_sec", 600)
            profile.audio_quality = h.get("audio_quality", 0.70)
            profile.chart_density = h.get("chart_density", 0.20)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "spreadsheet":
            profile.spreadsheet_complexity = h.get("spreadsheet_complexity",
                0.75 if any(k in p for k in ["model", "forecast", "multi-she"]) else 0.40)
            profile.table_density = h.get("table_density", 0.90)
            profile.detected_language = h.get("detected_language", "en")
        elif fmt == "presentation":
            profile.presentation_complexity = h.get("presentation_complexity", 0.50)
            profile.chart_density = h.get("chart_density", 0.55)
            profile.detected_language = h.get("detected_language", "en")
        else:
            profile.detected_language = h.get("detected_language", "en")
        profiles.append(profile)
    return profiles


PROMPT_MODALITY_HINTS: Dict[str, List[str]] = {
    "image": ["this image", "this picture", "this photo", "the photo"],
    "pdf": ["this pdf", "this document", "the document", "this contract"],
    "audio": ["this recording", "this audio", "this call", "the meeting recording"],
    "video": ["this video", "this clip", "the footage"],
    "spreadsheet": ["this spreadsheet", "this excel", "this workbook"],
    "presentation": ["this slide deck", "this presentation", "these slides"],
}

def detect_conflicts(prompt: str, profiles: List[ArtifactProfile]) -> List[str]:
    """Flag mismatches between the modality the prompt implies and what was uploaded.

    Resolution policy: trust the uploaded artifact, log a warning flag. The router
    keeps routing on real artifacts; the flag surfaces in the audit record so a
    reviewer (or upstream UX) can prompt the user if desired.
    """
    flags: List[str] = []
    p = prompt.lower()
    present_formats = {prof.format for prof in profiles}
    for modality, phrases in PROMPT_MODALITY_HINTS.items():
        if any(phrase in p for phrase in phrases) and modality not in present_formats:
            uploaded = sorted(present_formats) or ["none"]
            flags.append(
                f"prompt_implies_{modality}_but_uploaded_{'/'.join(uploaded)}"
            )
    return flags


TOPIC_KEYWORDS: Dict[str, List[str]] = {
    "legal_contract": ["agreement", "obligations", "liability", "indemnif",
                       "nda", "confidential", "party of the", "hereinafter", "clause"],
    "financial_document": ["invoice", "balance sheet", "revenue", "ebitda",
                           "statement", "tax", "amount due", "fiscal"],
    "medical_record": ["patient", "diagnosis", "treatment", "prescription",
                       "symptom", "clinical"],
    "research_paper": ["abstract", "we propose", "related work", "experiment",
                       "benchmark", "et al", "references"],
    "security_report": ["vulnerability", "cve", "exploit", "payload", "attack surface"],
    "support_record": ["ticket", "customer", "refund", "complaint", "resolution"],
}

def extract_text_preview(path: str, fmt: str, max_words: int = 500) -> Optional[str]:
    """Pull a small text preview from a file for topic inference. Best-effort."""
    try:
        if fmt == "pdf" and fitz is not None:
            doc = fitz.open(path)
            text = doc.load_page(0).get_text("text") or ""
            doc.close()
        elif fmt == "spreadsheet" and openpyxl is not None:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            ws = wb.active
            cells = []
            for row in ws.iter_rows(max_row=20, values_only=True):
                cells.extend(str(c) for c in row if c is not None)
            text = " ".join(cells)
            wb.close()
        elif fmt == "presentation" and pptx_mod is not None:
            from pptx import Presentation
            prs = Presentation(path)
            chunks = []
            for slide in list(prs.slides)[:3]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
            text = " ".join(chunks)
        elif fmt == "text":
            with open(path, "r", errors="ignore") as fh:
                text = fh.read(8000)
        else:
            return None
        return " ".join(text.split()[:max_words]) or None
    except Exception:
        return None

def infer_topic(preview: Optional[str]) -> Optional[str]:
    """Score the preview against TOPIC_KEYWORDS; return best match or None."""
    if not preview:
        return None
    low = preview.lower()
    best_topic, best_score = None, 0
    for topic, kws in TOPIC_KEYWORDS.items():
        score = sum(1 for kw in kws if kw in low)
        if score > best_score:
            best_topic, best_score = topic, score
    return best_topic if best_score >= 2 else None


def deterministic_extract(prompt: str, input_formats: List[str], estimated_tokens: int) -> Dict[str, Any]:
    p = prompt.lower()
    required_formats = sorted(set(input_formats) | {"text"})
    
    if any(k in p for k in ["image of", "picture of", "photo of", "generate an image", "draw", "render", "3d model", "cad", "obj file", "stl"]):
        if "image" not in required_formats: required_formats.append("image")
    if any(k in p for k in ["video of", "generate a video", "create a video", "animation", "animate"]):
        if "video" not in required_formats: required_formats.append("video")
    if any(k in p for k in ["generate music", "compose a song", "text to speech", "voiceover", "audiobook"]):
        if "audio" not in required_formats: required_formats.append("audio")
        
    return {
        "required_formats": required_formats,
        "requires_json": any(k in p for k in ["return json", "respond in json", "json output", "json schema", "structured json"]),
        "requires_function_calling": any(k in p for k in ["tool", "function call", "call a tool", "api call", "use tools"]),
        "requires_web_search": any(k in p for k in ["latest", "today", "current", "search the web", "news", "real-time"]),
        "requires_ocr": (
            any(fmt in input_formats for fmt in ["pdf", "image"])
            and any(k in p for k in ["scan", "extract text", "read the document", "ocr", "handwritten", "digitize"])
        ),
        "requires_citations": any(k in p for k in ["cite", "citations", "sources", "references", "bibliography"]),
        "requires_verifier": any(k in p for k in ["verify", "double-check", "validator", "review", "compliance", "audit"]),
        "min_context_window": estimated_tokens
    }


def is_simple_conversational_prompt(prompt: str) -> bool:
    import re
    normalized = re.sub(r'[^\w\s]', '', (prompt or "").strip().lower())
    words = [w for w in normalized.split() if w]
    simple_phrases = {
        "hi", "hello", "hey", "yo", "how are you", "hi how are you",
        "hello how are you", "hey how are you", "good morning", "good afternoon",
        "good evening", "thanks", "thank you"
    }
    is_simple = len(words) <= 6 and (normalized in simple_phrases or (words and words[0] in {"hi", "hello", "hey"}))
    return is_simple


def fallback_structured_semantic_parse(
    prompt: str,
    input_formats: List[str],
    estimated_tokens: int,
    artifacts: List[ArtifactProfile]
) -> StructuredSemanticParse:
    """Heuristic parser with intent scoring for multi-intent detection."""
    from collections import defaultdict
    p = prompt.lower()
    
    # 1. Keyword to Intent Mapping with weights
    intent_signals = {
        "audio": [("summarize", 0.5), ("chat", 0.5)],
        "ocr": [("extract text", 1.0), ("ocr", 1.0), ("scan", 0.8), ("handwritten", 1.0)],
        "document_qa": [("summarize", 0.8), ("contract", 1.0), ("document", 0.6), ("obligations", 0.8), ("clause", 1.0), ("review", 0.7), ("tabulate", 0.8)],
        "coding": [("python", 1.0), ("bug", 0.8), ("debug", 1.0), ("algorithm", 0.8), ("implement", 0.7), ("code", 1.0), ("sql", 1.0), ("function", 0.6), ("class", 0.5), ("website", 0.8), ("html", 1.0), ("css", 1.0), ("javascript", 1.0), ("react", 1.0), ("frontend", 0.8), ("backend", 0.8), ("script", 0.7), ("ui/ux", 1.0), ("app", 0.5)],
        "agent": [("automate", 1.0), ("scrape", 1.0), ("bot", 1.0), ("agent", 1.0), ("orchestrate", 1.0), ("multi-step plan", 1.0), ("workflow", 0.8), ("autonomous", 1.0)],
        "mathematics": [("prove", 0.8), ("integral", 1.0), ("equation", 1.0), ("theorem", 1.0), ("calculate", 0.6), ("derivative", 1.0), ("math", 1.0)],
        "translation": [("translate", 1.0), ("spanish", 0.8), ("french", 0.8), ("language", 0.5), ("localization", 1.0), ("german", 0.8)],
        "reasoning": [("analyze", 0.8), ("derive", 0.8), ("compare", 0.7), ("reason", 0.8), ("evaluate", 0.7), ("assess", 0.7), ("think", 0.5), ("logic", 0.8)],
        "creative": [("image of", 1.0), ("picture of", 1.0), ("photo of", 1.0), ("generate an image", 1.0), ("draw", 1.0), ("video of", 1.0), ("generate a video", 1.0), ("create a video", 1.0), ("animation", 1.0), ("animate", 1.0), ("3d model", 1.0), ("cad", 1.0), ("obj file", 1.0), ("stl", 1.0), ("render", 0.8), ("generate music", 1.0), ("compose a song", 1.0), ("text to speech", 1.0), ("voiceover", 1.0), ("audiobook", 1.0), ("roleplay", 1.0), ("pretend to be", 1.0), ("game", 0.7), ("text adventure", 1.0), ("dnd", 1.0), ("write a blog", 1.0), ("essay", 1.0), ("marketing", 0.8), ("copywriting", 1.0), ("story", 1.0), ("generate an article", 1.0), ("email template", 1.0), ("newsletter", 1.0), ("cover letter", 1.0), ("resume", 1.0), ("poem", 1.0)],
        "research": [("search", 0.8), ("research", 1.0), ("find information", 1.0), ("latest news", 1.0), ("current events", 1.0), ("who is", 0.8), ("what happened", 0.8)],
        "education": [("explain", 0.8), ("teach me", 1.0), ("how does", 0.8), ("tutorial", 1.0), ("learn", 0.8), ("for a 5 year old", 1.0), ("what is", 0.6)],
        "data_analysis": [("data", 0.6), ("csv", 1.0), ("analytics", 1.0), ("metrics", 1.0), ("chart", 1.0), ("graph", 1.0)],
        "summarization": [("tl;dr", 1.0), ("tldr", 1.0), ("summarize this", 1.0), ("bullet points", 0.8), ("key takeaways", 1.0), ("summarize", 0.7)]
    }

    # 2. Score Intents
    intent_scores = defaultdict(float)
    
    # Implicit format intents
    if any(fmt == "audio" for fmt in input_formats):
        intent_scores["audio"] += 2.0
    if any(fmt == "video" for fmt in input_formats):
        intent_scores["vision"] += 2.0
    if any(fmt == "pdf" for fmt in input_formats) or any(fmt == "spreadsheet" for fmt in input_formats):
        intent_scores["document_qa"] += 1.5
    if any(fmt == "image" for fmt in input_formats):
        if any(k in p for k in ["extract text", "ocr", "scan", "handwritten"]):
            intent_scores["ocr"] += 2.0
        else:
            intent_scores["vision"] += 1.0

    # Keyword scoring
    for intent, keywords in intent_signals.items():
        for kw, weight in keywords:
            if kw in p:
                intent_scores[intent] += weight

    # 3. Determine Primary and Secondary
    primary = "chat"
    secondary = []
    reason_parts = []
    
    if intent_scores:
        sorted_intents = sorted(intent_scores.items(), key=lambda x: x[1], reverse=True)
        primary = sorted_intents[0][0]
        secondary = [intent for intent, score in sorted_intents[1:] if score > 0.5]
        reason_parts.append(f"Intent scoring primary: {primary} ({sorted_intents[0][1]:.1f}).")
        if secondary:
            reason_parts.append(f"Detected secondary intents: {', '.join(secondary)}.")
    else:
        if is_simple_conversational_prompt(prompt):
            reason_parts.append("Simple conversational prompt detected.")
        else:
            reason_parts.append("No strong intents detected, defaulting to chat.")

    # 4. Domain, risk tier, risk type detection (scoring based as well)
    domain_signals = {
        "medical": ["medical", "diagnosis", "patient", "treatment", "symptom", "chest pain", "clinical", "prescription", "side effects", "disease", "therapy", "mental health", "depressed", "anxious", "suicidal", "need to talk", "lonely"],
        "legal": ["legal", "contract", "compliance", "law", "litigation", "obligation", "clause", "sue", "lawsuit", "court", "attorney", "nda", "terms of service"],
        "finance": ["finance", "tax", "investment", "trading", "portfolio", "stock", "crypto", "bitcoin", "mortgage", "loan", "interest rate"],
        "security": ["security", "vulnerability", "exploit", "incident", "breach", "ddos", "phishing", "malware", "ransomware", "bypass", "hack"],
        "sensitive": ["trump", "biden", "epstein", "election", "politics", "controversy", "nsfw", "porn", "violence", "kill", "bomb", "weapon", "illegal", "suicide", "self-harm", "drugs", "erotica", "sex", "nude", "naked", "fetish", "murder", "terrorist"],
        "customer_support": ["support ticket", "customer", "refund", "support"]
    }
    
    domain_scores = defaultdict(float)
    for dom, kws in domain_signals.items():
        for kw in kws:
            if kw in p:
                domain_scores[dom] += 1.0
                
    document_type = "generic"
    if domain_scores:
        top_domain = sorted(domain_scores.items(), key=lambda x: x[1], reverse=True)[0][0]
        if top_domain == "medical":
            domain, risk_tier, risk_type, document_type = "medical", "high", "regulated_advice", "medical_record"
        elif top_domain == "legal":
            domain, risk_tier, risk_type, document_type = "legal", "high", "regulated_advice", "contract"
        elif top_domain == "finance":
            domain, risk_tier, risk_type, document_type = "finance", "high", "regulated_advice", "financial_document"
        elif top_domain == "security":
            domain, risk_tier, risk_type, document_type = "security", "high", "security_sensitive", "security_report"
        elif top_domain == "sensitive":
            domain, risk_tier, risk_type, document_type = "sensitive", "high", "safety_sensitive", "controversial"
        elif top_domain == "customer_support":
            domain, risk_tier, risk_type, document_type = "customer_support", "low", "standard", "support_record"
    elif primary in ("coding", "agent"):
        domain, risk_tier, risk_type = "software", "medium", "operational"
    elif primary in ("reasoning", "mathematics"):
        domain, risk_tier, risk_type = "science", "medium", "analytical"
    else:
        domain, risk_tier, risk_type = "general", "low", "standard"
        
    # Expected output
    if any(k in p for k in ["json", "schema", "structured"]):
        expected_output = "structured_json"
    elif primary == "coding":
        expected_output = "code"
    else:
        expected_output = "free_text"
        
    # Ambiguity
    ambiguity_score = 0.15
    if any(k in p for k in ["maybe", "not sure", "either", "somehow", "approximately", "i think"]):
        ambiguity_score = 0.55
    if any(k in p for k in ["unclear", "ambiguous", "open-ended"]):
        ambiguity_score = 0.75
        
    if is_simple_conversational_prompt(prompt):
        ambiguity_score = 0.05
        
    # Actionability
    actionability = "advisory"
    if any(k in p for k in ["do this", "execute", "send", "file", "submit", "trade", "prescribe"]):
        actionability = "high"
        
    # Decomposition and stages
    required_stages = []
    decomposition_needed = False
    needs_verification = False
    
    combined_intents = [primary] + secondary
    
    if "ocr" in combined_intents:
        required_stages.append("ocr")
    if "vision" in combined_intents:
        required_stages.append("vision_understanding")
    if "audio" in combined_intents:
        required_stages.append("audio_understanding")
        
    if "document_qa" in combined_intents:
        required_stages.append("document_qa")
        if domain in {"legal", "medical", "finance"}:
            required_stages.append("domain_reasoning")
            decomposition_needed = True
            needs_verification = True
            
    for i in combined_intents:
        if i not in ["ocr", "vision", "audio", "document_qa"] and i not in required_stages:
            required_stages.append(i)
            
    if expected_output == "structured_json":
        required_stages.append("structured_output")
        
    if len(required_stages) == 0:
        required_stages.append(primary)
    elif len(required_stages) > 1:
        decomposition_needed = True
        
    # Build workflow graph
    workflow_graph = []
    for i, stage in enumerate(required_stages):
        workflow_graph.append({
            "stage_id": i + 1,
            "stage_name": stage,
            "depends_on": [] if i == 0 else [i]
        })
        
    # Parser confidence
    parser_confidence = 0.78
    if ambiguity_score > 0.5:
        parser_confidence -= 0.15
    if decomposition_needed:
        parser_confidence -= 0.05
        
    return StructuredSemanticParse(
        primary_family=primary,
        secondary_families=secondary,
        required_stages=required_stages,
        workflow_graph=workflow_graph,
        domain=domain,
        risk_tier=risk_tier,
        risk_type=risk_type,
        expected_output=expected_output,
        ambiguity_score=ambiguity_score,
        actionability=actionability,
        document_type=document_type,
        decomposition_needed=decomposition_needed,
        needs_verification=needs_verification,
        parser_confidence=parser_confidence,
        reason_summary=" ".join(reason_parts)
    )


def validate_structured_parse(data: StructuredSemanticParse) -> StructuredSemanticParse:
    allowed_families = {
        "coding", "reasoning", "mathematics", "chat", "vision",
        "ocr", "document_qa", "summarization", "translation", "agent", "audio",
        "creative", "research", "education", "data_analysis"
    }
    if data.primary_family not in allowed_families:
        raise ValueError(f"Invalid primary family: {data.primary_family}")
    if data.risk_tier not in {"low", "medium", "high"}:
        raise ValueError(f"Invalid risk tier: {data.risk_tier}")
    if data.expected_output not in {"structured_json", "free_text", "code"}:
        raise ValueError(f"Invalid expected output: {data.expected_output}")
    if not (0.0 <= data.ambiguity_score <= 1.0):
        raise ValueError("Ambiguity score must be in [0, 1].")
    return data


def parse_with_llm_interface(
    prompt: str,
    input_formats: List[str],
    estimated_tokens: int,
    artifacts: List[ArtifactProfile],
    llm_parser: Optional[Callable] = None
) -> StructuredSemanticParse:
    """Use real LLM parser if provided, else fall back to heuristic."""
    if llm_parser is not None:
        parsed = llm_parser(
            prompt=prompt,
            input_formats=input_formats,
            estimated_tokens=estimated_tokens,
            artifacts=[asdict(a) for a in artifacts]
        )
        if isinstance(parsed, StructuredSemanticParse):
            return validate_structured_parse(parsed)
        if isinstance(parsed, dict):
            return validate_structured_parse(StructuredSemanticParse(**parsed))
        raise ValueError("llm_parser returned unsupported type.")
    return validate_structured_parse(
        fallback_structured_semantic_parse(prompt, input_formats, estimated_tokens, artifacts)
    )


def infer_workflow_profile(
    primary_family: str,
    domain: str,
    input_formats: List[str],
    prompt: str,
    complexity: str = "low",
    risk_tier: str = "low",
) -> str:
    """Map task signals to an optimal weight profile.

    Priority: quality > riskfit > reliability > uncertainty > latency > runtime > cost
    High-risk and complex tasks always escalate to quality-dominant profiles.
    """
    p = prompt.lower()
    # ── Hard escalations: risk and complexity override everything ─────────────
    if risk_tier == "high":
        if domain == "legal" and "contract" in p:
            return "contract_review_intake"
        return "high_risk"
    if complexity == "high":
        return "complex_reasoning"
    if primary_family == "chat" and domain == "general" and is_simple_conversational_prompt(prompt):
        return "latency_first"
    # ── Domain-based quality profiles ─────────────────────────────────────────
    if domain in {"legal", "medical", "finance", "security"}:
        if domain == "legal" and "contract" in p:
            return "contract_review_intake"
        return "quality_first"
    if complexity == "medium":
        return "research_drafting"
    # ── Task-family profiles ───────────────────────────────────────────────────
    if primary_family == "coding":
        return "coding_assistant"
    if primary_family in ("creative", "research", "education"):
        return "quality_first"
    if primary_family == "data_analysis":
        return "research_drafting"
    if primary_family == "summarization":
        return "research_drafting"
    if primary_family == "ocr" and any(fmt in input_formats for fmt in ["image", "pdf"]):
        return "invoice_ocr_pipeline"
    if primary_family == "audio":
        return "real_time_voice_agent" if "real-time" in p else "audio_summary"
    if primary_family == "translation":
        return "multilingual_chat"
    if domain == "research":
        return "research_drafting"
    if "support" in p or domain == "customer_support":
        return "customer_support_summarization"
    return "generic_balanced"


def parse_task_request(
    prompt: str,
    input_formats: Optional[List[str]] = None,
    estimated_tokens: int = 2000,
    estimated_output_tokens: Optional[int] = None,
    artifact_hints: Optional[List[Dict[str, Any]]] = None,
    llm_parser: Optional[Callable] = None,
    request_constraints: Optional[RequestConstraints] = None,
    tenant_context: Optional[TenantContext] = None,
    files: Optional[List[str]] = None,
) -> TaskFeatures:
    rc = request_constraints or RequestConstraints()
    tc = tenant_context or TenantContext()

    # ---- Determine artifacts and formats ----
    if files:
        artifacts = inspect_artifacts(prompt=prompt, artifact_hints=artifact_hints, files=files)
        input_formats = sorted({a.format for a in artifacts} | {"text"})
    else:
        input_formats = input_formats or ["text"]
        bad = [f for f in input_formats if f not in SUPPORTED_FORMATS]
        if bad:
            raise ValueError(f"Unsupported formats: {bad}. Allowed: {SUPPORTED_FORMATS}")
        artifacts = inspect_artifacts(input_formats=input_formats, prompt=prompt,
                                      artifact_hints=artifact_hints)

    # Validate auto-detected formats too.
    bad = [f for f in input_formats if f not in SUPPORTED_FORMATS]
    if bad:
        raise ValueError(f"Unsupported formats: {bad}. Allowed: {SUPPORTED_FORMATS}")

    # ---- Conflict detection + semantic aggregation ----
    conflict_flags = detect_conflicts(prompt, artifacts)
    previews = [a.extracted_text_preview for a in artifacts if a.extracted_text_preview]
    extracted_text_summary = " ".join(previews)[:2000] if previews else None
    detected_languages = sorted({a.detected_language for a in artifacts if a.detected_language})
    total_file_size = sum(a.file_size_bytes or 0 for a in artifacts)

    hard = deterministic_extract(prompt, input_formats, estimated_tokens)
    # Artifact-driven OCR: a high scan_likelihood means OCR is needed even if
    # the prompt never says "scan" or "extract text".
    if not hard["requires_ocr"]:
        for a in artifacts:
            if a.format in {"pdf", "image"} and (a.scan_likelihood or 0) >= 0.6:
                hard["requires_ocr"] = True
                break

    soft = parse_with_llm_interface(prompt, input_formats, estimated_tokens, artifacts, llm_parser)

    # ---- Topic-driven domain/document override (semantic extraction wins when
    #      the prompt was uninformative, e.g. just "Explain this"). ----
    inferred_topics = [a.inferred_topic for a in artifacts if a.inferred_topic]
    if inferred_topics and soft.domain == "general":
        topic_to_domain = {
            "legal_contract": ("legal", "high", "regulated_advice", "contract"),
            "financial_document": ("finance", "high", "regulated_advice", "financial_document"),
            "medical_record": ("medical", "high", "regulated_advice", "medical_record"),
            "security_report": ("security", "high", "security_sensitive", "security_report"),
            "research_paper": ("research", "low", "standard", "research_paper"),
            "support_record": ("customer_support", "low", "standard", "support_record"),
        }
        topic = inferred_topics[0]
        if topic in topic_to_domain:
            soft.domain, soft.risk_tier, soft.risk_type, soft.document_type = topic_to_domain[topic]

    # Composite complexity
    comp_score = 0
    if estimated_tokens > 100000: comp_score += 2
    elif estimated_tokens > 20000: comp_score += 1
    if "reasoning" in soft.primary_family or "mathematics" in soft.primary_family: comp_score += 1
    if soft.risk_tier in ["high", "critical"]: comp_score += 1
    if len(soft.required_stages) > 1: comp_score += 1
    
    complexity = "high" if comp_score >= 3 else "medium" if comp_score >= 1 else "low"
    workflow_profile = infer_workflow_profile(soft.primary_family, soft.domain, input_formats, prompt, complexity, soft.risk_tier)

    requires_verifier = (
        hard["requires_verifier"]
        or soft.needs_verification
        or soft.risk_tier == "high"
        or rc.mandatory_verifier
    )
    safety_sensitive = soft.risk_tier == "high" or soft.risk_type in {"regulated_advice", "security_sensitive"}

    if rc.no_web_access:
        hard["requires_web_search"] = False

    return TaskFeatures(
        raw_prompt=prompt,
        input_formats=sorted(set(input_formats) | {"text"}),
        estimated_tokens=estimated_tokens,
        estimated_output_tokens=estimated_output_tokens,
        artifacts=artifacts,
        primary_family=soft.primary_family,
        secondary_families=soft.secondary_families,
        required_stages=soft.required_stages,
        workflow_graph=soft.workflow_graph,
        complexity=complexity,
        domain=soft.domain,
        risk_tier=soft.risk_tier,
        risk_type=soft.risk_type,
        required_formats=hard["required_formats"],
        requires_json=hard["requires_json"],
        requires_function_calling=hard["requires_function_calling"],
        requires_web_search=hard["requires_web_search"],
        requires_ocr=hard["requires_ocr"],
        requires_citations=hard["requires_citations"],
        requires_verifier=requires_verifier,
        min_context_window=hard["min_context_window"],
        expected_output=soft.expected_output,
        ambiguity_score=soft.ambiguity_score,
        safety_sensitive=safety_sensitive,
        actionability=soft.actionability,
        document_type=soft.document_type,
        decomposition_needed=soft.decomposition_needed,
        workflow_profile=workflow_profile,
        parser_confidence=soft.parser_confidence,
        conflict_flags=conflict_flags,
        extracted_text_summary=extracted_text_summary,
        detected_languages=detected_languages,
        total_file_size_bytes=total_file_size,
        request_constraints=rc,
        tenant_context=tc,
    )


def signal_value(model: Dict[str, Any], signal: str) -> float:
    try:
        kind, key = signal.split(":", 1)
    except ValueError:
        return 0.5
    if kind == "perf":
        return model.get("performance", {}).get(key, 0.5)
    if kind == "bench":
        val = model.get("benchmarks", {}).get(key, 0.5)
        return val / 100.0 if val > 1.0 else val
    if kind == "domain":
        return min(1.0, model.get("domains", {}).get(key, 0.5))
    return 0.5


def family_fit(model: Dict[str, Any], task: TaskFeatures) -> float:
    sigmap = TASK_FAMILY_SIGNAL_MAP.get(task.primary_family)
    if not sigmap:
        return model["domains"].get("general", 0.5)
    base = sum(signal_value(model, sig) * w for sig, w in sigmap.items())
    domkey = DOMAIN_EXPERTISE_MAP.get(task.domain, "general")
    domfit = min(1.0, model["domains"].get(domkey, 0.5))
    fit = 0.75 * base + 0.25 * domfit
    if task.expected_output == "structured_json":
        fit = 0.70 * fit + 0.30 * model["performance"].get("json_reliability", 0.5)
    if task.min_context_window > 100_000:
        fit = 0.70 * fit + 0.30 * model["performance"].get("long_context", 0.5)
    if task.requires_ocr:
        fit = 0.70 * fit + 0.30 * model["performance"].get("ocr", 0.5)
    if "spreadsheet" in task.input_formats:
        fit = 0.80 * fit + 0.20 * model["performance"].get("spreadsheet_reasoning", 0.5)
    if "audio" in task.input_formats:
        fit = 0.70 * fit + 0.30 * model["performance"].get("audio_understanding", 0.5)
    return float(min(max(fit, 0.0), 1.0))


def stage_fit(model: Dict[str, Any], stage_name: str, task: TaskFeatures) -> float:
    stage_map = {
        "domain_reasoning": lambda: min(1.0, model["domains"].get(DOMAIN_EXPERTISE_MAP.get(task.domain, "general"), 0.5)),
        "structured_output": lambda: model["performance"].get("json_reliability", 0.5),
        "audio_understanding": lambda: model["performance"].get("audio_understanding", 0.5),
        "vision_understanding": lambda: model["performance"].get("vision_understanding", 0.5),
        "ocr": lambda: model["performance"].get("ocr", 0.5),
        "coding": lambda: model["performance"].get("coding", 0.5),
        "document_qa": lambda: model["performance"].get("document_qa", 0.5),
        "summarization": lambda: model["performance"].get("summarization", 0.5),
    }
    if stage_name in stage_map:
        return stage_map[stage_name]()
    return family_fit(model, task)


def composite_workflow_fit(model: Dict[str, Any], task: TaskFeatures) -> float:
    scores = [stage_fit(model, stage, task) for stage in task.required_stages]
    if not scores:
        return family_fit(model, task)
    return float(statistics.mean(scores))


def risk_support(model: Dict[str, Any], task: TaskFeatures) -> float:
    safety_key = DOMAIN_SAFETY_MAP.get(task.domain, "general")
    base = model["safety"].get(safety_key, model["safety"].get("general", 0.7))
    if task.risk_type == "regulated_advice":
        base = 0.8 * base + 0.2 * model["safety"].get("regulated_advice", base)
    if task.risk_tier == "low":
        return min(1.0, base + 0.05)
    return base


def static_cost_score(model: Dict[str, Any]) -> float:
    return model["pricing"]["relative_cost_score"]


def static_latency_score(model: Dict[str, Any]) -> float:
    return model["ops_static"]["latency_score"]


def reliability_score(model: Dict[str, Any]) -> float:
    return model["ops_static"]["reliability"]


def freshness_score(model: Dict[str, Any]) -> float:
    sec = model["ops_dynamic"].get("telemetry_freshness_sec", 60)
    return max(0.0, 1.0 - min(sec / 600.0, 1.0))


def incident_penalty(status: str) -> float:
    return {"green": 0.0, "yellow": 0.08, "orange": 0.18, "red": 0.35}.get(status, 0.10)


def runtime_health_score(model: Dict[str, Any]) -> float:
    dyn = model["ops_dynamic"]
    components = [
        (0.20, max(0.0, 1.0 - min(dyn["recent_latency_ms"] / 5000.0, 1.0))),
        (0.18, max(0.0, 1.0 - min(dyn["recent_failure_rate"] / 0.15, 1.0))),
        (0.16, dyn["current_availability"]),
        (0.10, max(0.0, 1.0 - dyn["rate_limit_pressure"])),
        (0.08, max(0.0, 1.0 - dyn["queue_pressure"])),
        (0.10, max(0.0, 1.0 - incident_penalty(dyn["incident_status"]))),
        (0.08, max(0.0, 1.0 - 0.5 * dyn["budget_pressure"])),
        (0.10, freshness_score(model)),
    ]
    return float(min(max(sum(w * v for w, v in components), 0.0), 1.0))


def effective_prior(model: Dict[str, Any], family: str) -> Tuple[float, float]:
    fam = model.get("priors", {}).get("task_family", {}).get(family, model.get("priors", {}).get("global", {"alpha": 1.0, "beta": 1.0}))
    alpha, beta = float(fam["alpha"]), float(fam["beta"])
    ev = model.get("evaluation", {})
    if ev.get("samples", 0) > 0:
        alpha += ev.get("wins", 0)
        beta += ev.get("losses", 0)
    return alpha, beta


def estimate_request_cost_usd(model: Dict[str, Any], input_tokens: Optional[int], output_tokens: Optional[int], is_cached: bool = False) -> float:
    if input_tokens is None:
        input_tokens = 500
    if output_tokens is None:
        output_tokens = 50  # Provide a safe default for cost estimation when not specified
    return ((input_tokens / 1_000_000) * model["pricing"]["input_cost"] +
            (output_tokens / 1_000_000) * model["pricing"]["output_cost"])


def estimate_request_latency_ms(model: Dict[str, Any], task: TaskFeatures, n_stages: int = 1) -> float:
    lp = model.get("latency_profile", {})
    ttft = lp.get("ttft_ms_mean", model.get("ops_dynamic", {}).get("recent_latency_ms", 500) * 0.3)
    tpot = lp.get("tpot_ms_mean", 15)
    out_tokens = getattr(task, 'estimated_output_tokens', None)
    if out_tokens is None:
        out_tokens = 1000 if task.complexity == "high" else (250 if task.complexity == "medium" else 50)
    base = ttft + (out_tokens * tpot)
    complexity_mult = {"low": 1.0, "medium": 1.25, "high": 1.60}.get(task.complexity, 1.0)
    modality_mult = 1.0
    if "audio" in task.input_formats or "video" in task.input_formats:
        modality_mult = 1.35
    elif "image" in task.input_formats or "pdf" in task.input_formats:
        modality_mult = 1.15
    json_mult = 1.10 if task.requires_json else 1.0
    return float(base * complexity_mult * modality_mult * json_mult * n_stages)


def feasibility_filter(
    task: TaskFeatures,
    registry: List[Dict[str, Any]]
) -> Tuple[List[Dict[str, Any]], Dict[str, str]]:
    feasible = []
    reasons = {}
    rc = task.request_constraints
    tc = task.tenant_context
    for model in registry:
        name = model["name"]
        if model.get("status") != "active" or not model.get("api_available", False):
            reasons[name] = "inactive_or_no_api"
            continue
        if model["ops_dynamic"].get("incident_status") == "red":
            reasons[name] = "runtime_incident_red"
            continue
        if rc.allowed_providers and model["provider"] not in rc.allowed_providers:
            reasons[name] = "provider_not_allowed"
            continue
        if rc.disallowed_providers and model["provider"] in rc.disallowed_providers:
            reasons[name] = "provider_disallowed"
            continue
        if rc.allowed_tiers and model["tier"] not in rc.allowed_tiers:
            reasons[name] = "tier_not_allowed"
            continue
        if rc.no_open_weight and model.get("open_weight", False):
            reasons[name] = "open_weight_disallowed"
            continue
        if tc.allowed_models and name not in tc.allowed_models:
            reasons[name] = "tenant_model_not_allowed"
            continue
        if rc.required_region:
            regions = model.get("allowed_regions", [])
            if rc.required_region not in regions and "global" not in regions:
                reasons[name] = "region_not_supported"
                continue
        mods = model["modalities"]
        caps = model["capabilities"]
        ctx = model["context"]
        missing_fmt = [f for f in task.required_formats if not mods.get(f, False)]
        if missing_fmt:
            reasons[name] = f"missing_format:{','.join(missing_fmt)}"
            continue
        if task.requires_json and not caps.get("json_mode", False):
            reasons[name] = "missing_json_mode"
            continue
        if task.requires_function_calling and not caps.get("function_calling", False):
            reasons[name] = "missing_function_calling"
            continue
        if task.requires_web_search and not caps.get("web_search", False):
            reasons[name] = "missing_web_search"
            continue
        if task.requires_ocr and not caps.get("ocr", False):
            reasons[name] = "missing_ocr"
            continue
        if task.requires_citations and not caps.get("citation_support", False):
            reasons[name] = "missing_citation_support"
            continue
        if ctx["window"] < task.min_context_window:
            reasons[name] = "insufficient_context_window"
            continue
        # Budget pre-check
        est_cost = estimate_request_cost_usd(model, task.estimated_tokens, task.estimated_output_tokens, getattr(task, 'requires_context_caching', False))
        if rc.max_cost_usd is not None and est_cost > rc.max_cost_usd:
            reasons[name] = f"estimated_cost_exceeds_budget:{est_cost:.4f}"
            continue
        if tc.budget_remaining_usd is not None and est_cost > tc.budget_remaining_usd:
            reasons[name] = f"exceeds_tenant_remaining_budget:{est_cost:.4f}"
            continue
        feasible.append(deepcopy(model))
    return feasible, reasons


def evaluate_policy(task: TaskFeatures) -> PolicyDecision:
    notes: List[str] = []
    restricted_to_tiers: List[str] = []
    restricted_to_models: List[str] = []
    restricted_to_providers: List[str] = []
    require_verifier_types: List[str] = []
    overlay = task.tenant_context.policy_overlay or {}
    # Medical high-actionability escalation
    if task.domain == "medical" and task.actionability == "high":
        notes.append("Medical high-actionability task requires mandatory escalation.")
        require_verifier_types.extend(["safety_review", "factuality_review"])
        return PolicyDecision(
            allowed=True, must_escalate=True, must_abstain=False,
            restricted_to_tiers=["Frontier"],
            require_verifier_types=sorted(set(require_verifier_types)),
            notes=notes
        )
    # Legal tasks
    if task.domain == "legal":
        notes.append("Legal task requires policy scrutiny and verification.")
        require_verifier_types.extend(["safety_review", "schema_validation"])
        restricted_to_tiers.append("Frontier")
    # Finance investment tasks
    if task.domain == "finance" and any(k in task.raw_prompt.lower() for k in ["investment", "trade", "portfolio"]):
        notes.append("Investment finance task restricted to frontier models.")
        restricted_to_tiers.append("Frontier")
        require_verifier_types.extend(["factuality_review", "safety_review"])
    # Sensitive content escalation
    if task.domain == "sensitive":
        notes.append("Sensitive content detected — restricting to frontier models with safety review.")
        restricted_to_tiers.append("Frontier")
        require_verifier_types.extend(["safety_review", "factuality_review"])
    # Offensive security denial
    if task.domain == "security" and any(k in task.raw_prompt.lower() for k in ["exploit", "offensive", "payload"]):
        return PolicyDecision(
            allowed=False, must_escalate=False, must_abstain=True,
            deny_reason="restricted_security_content",
            notes=["Offensive security request denied by policy."]
        )
    # General verifier requirements
    if task.requires_verifier:
        vtype = "schema_validation" if task.expected_output == "structured_json" else "factuality_review"
        require_verifier_types.append(vtype)
    if task.requires_citations:
        require_verifier_types.append("citation_review")
    # High risk tier
    if task.risk_tier == "high":
        notes.append("High-risk task restricted to frontier-grade models.")
        restricted_to_tiers.append("Frontier")
    # Tenant overlays
    if overlay.get("frontier_only", False):
        restricted_to_tiers.append("Frontier")
        notes.append("Tenant overlay enforces frontier-only routing.")
    if overlay.get("allowed_providers"):
        restricted_to_providers.extend(overlay["allowed_providers"])
        notes.append("Tenant overlay restricts providers.")
    if overlay.get("mandatory_verifier", False):
        require_verifier_types.append("consistency_review")
        notes.append("Tenant overlay requires mandatory verification.")
    return PolicyDecision(
        allowed=True, must_escalate=False, must_abstain=False,
        restricted_to_tiers=sorted(set(restricted_to_tiers)),
        restricted_to_models=sorted(set(restricted_to_models)),
        restricted_to_providers=sorted(set(restricted_to_providers)),
        require_verifier_types=sorted(set(require_verifier_types)),
        notes=notes
    )


def apply_policy_to_models(
    task: TaskFeatures,
    policy: PolicyDecision,
    models: List[Dict[str, Any]]
) -> Tuple[List[Dict[str, Any]], Dict[str, str]]:
    if not policy.allowed:
        return [], {m["name"]: policy.deny_reason or "policy_denied" for m in models}
    gated = []
    reasons = {}
    for model in models:
        name = model["name"]
        rs = risk_support(model, task)
        if task.risk_tier == "high" and rs < 0.85:
            reasons[name] = f"insufficient_safety:{rs:.2f}"
            continue
        if task.risk_tier == "high" and "high_risk" in model["routing"].get("avoid_for", []):
            reasons[name] = "routing_avoid_high_risk"
            continue
        if policy.restricted_to_tiers and not any(t in model["tier"] for t in policy.restricted_to_tiers):
            reasons[name] = f"policy_tier_restriction:{model['tier']}"
            continue
        if policy.restricted_to_models and name not in policy.restricted_to_models:
            reasons[name] = "not_in_policy_model_allowlist"
            continue
        if policy.restricted_to_providers and model["provider"] not in policy.restricted_to_providers:
            reasons[name] = "not_in_policy_provider_allowlist"
            continue
        gated.append(model)
    return gated, reasons


def beta_mean(alpha: float, beta: float) -> float:
    return alpha / (alpha + beta)


def beta_variance(alpha: float, beta: float) -> float:
    return (alpha * beta) / (((alpha + beta) ** 2) * (alpha + beta + 1))


def convex_combine(values: List[float], weights: List[float]) -> float:
    total = sum(weights)
    if total == 0:
        return float(statistics.mean(values))
    return sum(v * w for v, w in zip(values, weights)) / total


def attach_contextual_quality(
    models: List[Dict[str, Any]],
    task: TaskFeatures
) -> List[Dict[str, Any]]:
    enriched = []
    for model in models:
        model = deepcopy(model)
        g = model.get("priors", {}).get("global", {"alpha": 1.0, "beta": 1.0})
        global_mean = beta_mean(g["alpha"], g["beta"])
        falpha, fbeta = effective_prior(model, task.primary_family)
        fam_mean = beta_mean(falpha, fbeta)
        fam_var = beta_variance(falpha, fbeta)
        fam_fit = family_fit(model, task)
        wf_fit = composite_workflow_fit(model, task)
        dom_fit = min(1.0, model["domains"].get(DOMAIN_EXPERTISE_MAP.get(task.domain, "general"), 0.5))
        rt_fit = runtime_health_score(model)
        fr_fit = freshness_score(model)
        contextual_mean = convex_combine(
            [global_mean, fam_mean, fam_fit, wf_fit, dom_fit, fr_fit],
            [0.12, 0.22, 0.18, 0.23, 0.15, 0.10]
        )
        runtime_adjusted_mean = 0.76 * contextual_mean + 0.24 * rt_fit
        model["q"] = {
            "global_mean": global_mean,
            "family_mean": fam_mean,
            "family_variance": fam_var,
            "family_fit": fam_fit,
            "workflow_fit": wf_fit,
            "domain_fit": dom_fit,
            "runtime_fit": rt_fit,
            "freshness_fit": fr_fit,
            "contextual_mean": contextual_mean,
            "runtime_adjusted_mean": runtime_adjusted_mean,
            "uncertainty": math.sqrt(max(fam_var, 1e-9)),
            "alpha": falpha,
            "beta": fbeta,
        }
        enriched.append(model)
    return enriched


def pareto_vector(model: Dict[str, Any]) -> Tuple[float, float, float, float, float]:
    return (
        model["q"]["runtime_adjusted_mean"],
        static_cost_score(model),
        static_latency_score(model),
        reliability_score(model),
        runtime_health_score(model),
    )


def dominates_with_margin(m1: Dict[str, Any], m2: Dict[str, Any], eps: float = 0.02) -> bool:
    v1 = pareto_vector(m1)
    v2 = pareto_vector(m2)
    ge_all = all(a >= b - eps for a, b in zip(v1, v2))
    gt_any = any(a > b + eps for a, b in zip(v1, v2))
    return ge_all and gt_any


def pareto_frontier(models: List[Dict[str, Any]], eps: float = 0.02) -> List[Dict[str, Any]]:
    frontier = []
    for i, model in enumerate(models):
        dominated = any(
            dominates_with_margin(models[j], model, eps=eps)
            for j in range(len(models)) if j != i
        )
        if not dominated:
            frontier.append(model)
    return frontier


WEIGHT_PROFILES = {
    "quality_first": {"quality": 0.44, "uncertainty": 0.16, "cost": 0.03, "latency": 0.05, "reliability": 0.16, "riskfit": 0.1, "runtime": 0.06},
    "complex_reasoning": {"quality": 0.42, "uncertainty": 0.16, "cost": 0.02, "latency": 0.04, "reliability": 0.16, "riskfit": 0.14, "runtime": 0.06},
    "high_risk": {"quality": 0.36, "uncertainty": 0.18, "cost": 0.01, "latency": 0.03, "reliability": 0.2, "riskfit": 0.18, "runtime": 0.04},
    "contract_review_intake": {"quality": 0.38, "uncertainty": 0.16, "cost": 0.02, "latency": 0.03, "reliability": 0.18, "riskfit": 0.18, "runtime": 0.05},
    "research_drafting": {"quality": 0.38, "uncertainty": 0.14, "cost": 0.05, "latency": 0.05, "reliability": 0.16, "riskfit": 0.12, "runtime": 0.1},
    "coding_assistant": {"quality": 0.36, "uncertainty": 0.14, "cost": 0.08, "latency": 0.08, "reliability": 0.16, "riskfit": 0.08, "runtime": 0.1},
    "balanced": {"quality": 0.32, "uncertainty": 0.12, "cost": 0.12, "latency": 0.12, "reliability": 0.14, "riskfit": 0.1, "runtime": 0.08},
    "generic_balanced": {"quality": 0.32, "uncertainty": 0.12, "cost": 0.12, "latency": 0.12, "reliability": 0.14, "riskfit": 0.1, "runtime": 0.08},
    "invoice_ocr_pipeline": {"quality": 0.32, "uncertainty": 0.14, "cost": 0.08, "latency": 0.08, "reliability": 0.14, "riskfit": 0.12, "runtime": 0.12},
    "audio_summary": {"quality": 0.28, "uncertainty": 0.12, "cost": 0.1, "latency": 0.18, "reliability": 0.14, "riskfit": 0.08, "runtime": 0.1},
    "multilingual_chat": {"quality": 0.26, "uncertainty": 0.1, "cost": 0.12, "latency": 0.2, "reliability": 0.14, "riskfit": 0.08, "runtime": 0.1},
    "real_time_voice_agent": {"quality": 0.18, "uncertainty": 0.08, "cost": 0.08, "latency": 0.32, "reliability": 0.14, "riskfit": 0.06, "runtime": 0.14},
    "customer_support_summarization": {"quality": 0.24, "uncertainty": 0.1, "cost": 0.24, "latency": 0.16, "reliability": 0.12, "riskfit": 0.06, "runtime": 0.08},
    "budget_first": {"quality": 0.22, "uncertainty": 0.1, "cost": 0.32, "latency": 0.12, "reliability": 0.1, "riskfit": 0.06, "runtime": 0.08},
    "latency_first": {"quality": 0.18, "uncertainty": 0.08, "cost": 0.08, "latency": 0.36, "reliability": 0.12, "riskfit": 0.06, "runtime": 0.12},
}


def minmax_normalize(values: List[float], invert: bool = False) -> List[float]:
    if not values:
        return []
    vmin, vmax = min(values), max(values)
    if math.isclose(vmin, vmax):
        return [0.5] * len(values)
    base = [(v - vmin) / (vmax - vmin) for v in values]
    return [1 - x for x in base] if invert else base


def choose_effective_profile(task: TaskFeatures, requested: str) -> str:
    if task.risk_tier == "high" and task.workflow_profile != "contract_review_intake":
        return "high_risk"
    if task.workflow_profile in WEIGHT_PROFILES:
        return task.workflow_profile
    return requested if requested in WEIGHT_PROFILES else "balanced"


def compute_utilities(
    models: List[Dict[str, Any]],
    task: TaskFeatures,
    profile_name: str = "balanced"
) -> List[Dict[str, Any]]:
    profile_name = choose_effective_profile(task, profile_name)
    w = WEIGHT_PROFILES[profile_name]
    qvals = [m["q"]["runtime_adjusted_mean"] for m in models]
    uvals = [m["q"]["uncertainty"] for m in models]
    cvals = [static_cost_score(m) for m in models]
    lvals = [static_latency_score(m) for m in models]
    rvals = [reliability_score(m) for m in models]
    rfvals = [risk_support(m, task) for m in models]
    rtvals = [runtime_health_score(m) for m in models]
    qn = minmax_normalize(qvals)
    un = minmax_normalize(uvals)
    cn = minmax_normalize(cvals)
    ln = minmax_normalize(lvals)
    rn = minmax_normalize(rvals)
    rfn = minmax_normalize(rfvals)
    rtn = minmax_normalize(rtvals)
    out = []
    # Complexity modifier: dynamically shift weights at runtime
    effective_w = dict(w)
    if task.complexity == "high":
        quality_bonus = 0.10
        cost_cut = min(effective_w.get("cost", 0), quality_bonus * 0.70)
        lat_cut  = min(effective_w.get("latency", 0), quality_bonus * 0.30)
        effective_w["quality"]  = min(0.55, effective_w.get("quality", 0) + quality_bonus)
        effective_w["cost"]     = max(0.01, effective_w.get("cost", 0)    - cost_cut)
        effective_w["latency"]  = max(0.02, effective_w.get("latency", 0) - lat_cut)
    elif task.complexity == "medium":
        quality_bonus = 0.05
        cost_cut = min(effective_w.get("cost", 0), quality_bonus)
        effective_w["quality"] = min(0.50, effective_w.get("quality", 0) + quality_bonus)
        effective_w["cost"]    = max(0.02, effective_w.get("cost", 0)    - cost_cut)
    elif task.complexity == "low":
        quality_cut = 0.15
        effective_w["quality"] = max(0.05, effective_w.get("quality", 0) - quality_cut)
        effective_w["cost"]    = min(0.50, effective_w.get("cost", 0) + (quality_cut * 0.6))
        effective_w["latency"] = min(0.50, effective_w.get("latency", 0) + (quality_cut * 0.4))
    w = effective_w

    for model, q, u, c, l, r, rf, rt in zip(models, qn, un, cn, ln, rn, rfn, rtn):
        m = deepcopy(model)
        # Capability gate: penalise weak models on complex/high-risk tasks
        capability_gate = 0.0
        if task.complexity == "high" or task.risk_tier == "high":
            human_eval = model.get("performance", {}).get("human_eval_score", 50)
            mMLU       = model.get("performance", {}).get("mMLU_score", 50)
            capability_gate = max(0.0, (70 - min(human_eval, mMLU)) / 100)
        utility = (
            w.get("quality", 0.0) * q
            - w.get("uncertainty", 0.0) * u
            + w.get("cost", 0.0) * c
            + w.get("latency", 0.0) * l
            + w.get("reliability", 0.0) * r
            + w.get("riskfit", 0.0) * rf
            + w.get("runtime", 0.0) * rt
            - capability_gate
        )
        # Penalty for ambiguous or highly decomposed tasks
        if task.ambiguity_score > 0.5:
            utility -= 0.03
        if task.decomposition_needed and len(task.required_stages) > 2:
            utility -= 0.01
        # SLA hard penalties
        rc = task.request_constraints
        est_latency = estimate_request_latency_ms(model, task)
        est_cost = estimate_request_cost_usd(model, task.estimated_tokens, task.estimated_output_tokens, getattr(task, 'requires_context_caching', False))
        sla_violation = False
        if rc.max_latency_ms is not None and est_latency > rc.max_latency_ms:
            utility -= 0.15
            sla_violation = True
        if rc.max_cost_usd is not None and est_cost > rc.max_cost_usd:
            utility -= 0.15
            sla_violation = True
        m["u"] = {
            "profile_used": profile_name,
            "expected_utility": utility,
            "n_quality": q, "n_uncertainty": u, "n_cost": c,
            "n_latency": l, "n_reliability": r, "n_riskfit": rf, "n_runtime": rt,
            "est_latency_ms": est_latency,
            "est_cost_usd": est_cost,
            "sla_violation": sla_violation,
        }
        out.append(m)
    return sorted(out, key=lambda x: x["u"]["expected_utility"], reverse=True)


def bounded_beta_sample(mean: float, variance: float) -> float:
    max_var = max(mean * (1 - mean) - 1e-6, 1e-6)
    variance = max(min(variance, max_var), 1e-6)
    common = (mean * (1 - mean) / variance) - 1
    alpha = max(mean * common, 1.0)
    beta_param = max((1 - mean) * common, 1.0)
    return float(GLOBAL_RNG.beta(alpha, beta_param))


def estimate_confidence(
    models: List[Dict[str, Any]],
    task: TaskFeatures,
    profile_name: str = "balanced",
    n_sim: int = 1500
) -> Dict[str, Any]:
    profile_name = choose_effective_profile(task, profile_name)
    w = WEIGHT_PROFILES[profile_name]
    wins = {m["name"]: 0 for m in models}
    for _ in range(n_sim):
        sampled_utilities = []
        for m in models:
            sq = bounded_beta_sample(m["q"]["runtime_adjusted_mean"], max(m["q"]["family_variance"], 1e-4))
            # Sampled utility aligned with the real utility's dominant terms
            su = (
                w.get("quality", 0.0) * sq
                + w.get("cost", 0.0) * static_cost_score(m)
                + w.get("latency", 0.0) * min(1.0, max(0.0, bounded_beta_sample(static_latency_score(m), 0.05)))
                + w.get("reliability", 0.0) * reliability_score(m)
                + w.get("runtime", 0.0) * runtime_health_score(m)
            )
            sampled_utilities.append((m["name"], su))
        winner = max(sampled_utilities, key=lambda x: x[1])[0]
        wins[winner] += 1
    win_probs = {k: v / n_sim for k, v in wins.items()}
    ordered = sorted(win_probs.items(), key=lambda x: x[1], reverse=True)
    top_name, top_prob = ordered[0]
    second_prob = ordered[1][1] if len(ordered) > 1 else 0.0
    return {
        "win_probabilities": win_probs,
        "top_model": top_name,
        "top_confidence": top_prob,
        "margin": top_prob - second_prob,
        "ordered": ordered,
        "n_simulations": n_sim,
        "calibration_note": "Calibrate against held-out labeled routing outcomes for production use."
    }


def choose_verifier_models(
    all_models: List[Dict[str, Any]],
    verifier_types: List[str],
    selected_model_name: Optional[str],
    max_verifiers: int = 2
) -> List[Dict[str, Any]]:
    """Select best verifier models for each required type, excluding the primary model."""
    if not verifier_types:
        return []
    candidates = [m for m in all_models if m["name"] != selected_model_name]
    if not candidates:
        return []
    scored = []
    for model in candidates:
        vf = model.get("verifier_fit", {})
        avg_fit = float(statistics.mean([min(1.0, vf.get(vt, 0.5)) for vt in verifier_types]))
        scored.append((model, avg_fit))
    scored.sort(key=lambda x: x[1], reverse=True)
    return [s[0] for s in scored[:max_verifiers]]


def select_best_for_stage(
    models: List[Dict[str, Any]],
    stage_name: str,
    task: TaskFeatures
) -> Tuple[Optional[Dict[str, Any]], List[Dict[str, Any]]]:
    if not models:
        return None, []
    scored = []
    for m in models:
        sf = stage_fit(m, stage_name, task)
        blended = 0.80 * sf + 0.20 * runtime_health_score(m)
        scored.append((m, blended))
    scored.sort(key=lambda x: x[1], reverse=True)
    best = scored[0][0]
    fallbacks = [s[0] for s in scored[1:3]]
    return best, fallbacks


def generate_multi_stage_plan(
    models: List[Dict[str, Any]],
    task: TaskFeatures,
    confidence_data: Dict[str, Any],
    verifier_models: List[Dict[str, Any]],
    policy: PolicyDecision,
    profile_name: str
) -> Optional[ExecutionPlan]:
    if len(task.required_stages) < 2:
        return None
    n_stages = len(task.required_stages)
    in_per_stage = max(task.estimated_tokens // n_stages, 1)
    out_per_stage = max((task.estimated_output_tokens or 500) // n_stages, 1)
    stage_routes: List[StageRoute] = []
    total_latency = 0.0
    total_cost = 0.0
    stage_quality_scores: List[float] = []
    for i, stage_name in enumerate(task.required_stages):
        best, fallbacks = select_best_for_stage(models, stage_name, task)
        if best is None:
            return None
        s_latency = estimate_request_latency_ms(best, task, n_stages=1) / max(n_stages, 1)
        s_cost = estimate_request_cost_usd(best, in_per_stage, out_per_stage)
        total_latency += s_latency
        total_cost += s_cost
        stage_quality_scores.append(stage_fit(best, stage_name, task))
        stage_routes.append(StageRoute(
            stage_id=i + 1,
            stage_name=stage_name,
            selected_model=best["name"],
            fallback_models=[fb["name"] for fb in fallbacks],
            verifier_models=[v["name"] for v in verifier_models] if i == n_stages - 1 else [],
            stage_confidence = confidence_data["win_probabilities"].get(best["name"], 0.0),
            expected_latency_ms=s_latency,
            expected_cost_usd=s_cost,
            explanation=f"Stage '{stage_name}' assigned to {best['name']} (best stage fit)."
        ))
    for v in verifier_models:
        total_latency += v["ops_dynamic"]["recent_latency_ms"] * 0.5
        total_cost += estimate_request_cost_usd(v, out_per_stage, 500)
    expected_quality = float(statistics.mean(stage_quality_scores)) if stage_quality_scores else 0.0
    return ExecutionPlan(
        plan_id=str(uuid.uuid4())[:12],
        plan_type="multi_stage",
        selected_model=None,
        stage_routes=stage_routes,
        fallback_models=[],
        verifier_models=[v["name"] for v in verifier_models],
        expected_latency_ms=total_latency,
        expected_cost_usd=total_cost,
        expected_quality=expected_quality,
        confidence = confidence_data["win_probabilities"].get(best["name"], 0.0),
        utility=0.0,
        confidence_margin=confidence_data["margin"],
        profile_used=profile_name,
        explanation={
            "plan_type": "multi_stage",
            "stage_assignments": [
                {"stage": sr.stage_name, "model": sr.selected_model} for sr in stage_routes
            ],
            "task_family": task.primary_family,
            "domain": task.domain,
            "risk_tier": task.risk_tier,
            "workflow_profile": task.workflow_profile,
            "verifiers_attached": [v["name"] for v in verifier_models],
            "policy_notes": policy.notes,
            "note": "Per-stage specialization; sequential latency is additive.",
        },
        trace={
            "router_version": ROUTER_VERSION,
            "scoring_version": SCORING_VERSION,
            "timestamp_utc": datetime.now(timezone.utc).isoformat(),
        }
    )


def generate_single_model_plan(
    model: Dict[str, Any],
    task: TaskFeatures,
    confidence_data: Dict[str, Any],
    verifier_models: List[Dict[str, Any]],
    policy: PolicyDecision,
    fallback_models: List[Dict[str, Any]],
    profile_name: str
) -> ExecutionPlan:
    est_latency = estimate_request_latency_ms(model, task)
    est_cost = estimate_request_cost_usd(model, task.estimated_tokens, task.estimated_output_tokens, getattr(task, 'requires_context_caching', False))
    stage_routes = []
    n_stages = max(len(task.required_stages), 1)
    for i, stage_name in enumerate(task.required_stages):
        stage_routes.append(StageRoute(
            stage_id=i + 1,
            stage_name=stage_name,
            selected_model=model["name"],
            fallback_models=[fb["name"] for fb in fallback_models[:2]],
            verifier_models=[v["name"] for v in verifier_models],
            stage_confidence = confidence_data["win_probabilities"].get(model["name"], 0.0),
            expected_latency_ms=est_latency / n_stages,
            expected_cost_usd=est_cost / n_stages,
            explanation=f"Stage '{stage_name}' handled by {model['name']}."
        ))
    # Add verifier overhead
    for v in verifier_models:
        v_latency = v["ops_dynamic"]["recent_latency_ms"] * 0.5
        v_cost = estimate_request_cost_usd(v, task.estimated_output_tokens, 500)
        est_latency += v_latency
        est_cost += v_cost
    return ExecutionPlan(
        plan_id=str(uuid.uuid4())[:12],
        plan_type="single_model",
        selected_model=model["name"],
        stage_routes=stage_routes,
        fallback_models=[fb["name"] for fb in fallback_models[:3]],
        verifier_models=[v["name"] for v in verifier_models],
        expected_latency_ms=est_latency,
        expected_cost_usd=est_cost,
        expected_quality=model["q"]["runtime_adjusted_mean"],
        confidence = confidence_data["win_probabilities"].get(model["name"], 0.0),
        utility=model["u"]["expected_utility"],
        confidence_margin=confidence_data["margin"],
        profile_used=profile_name,
        explanation={
            "primary_model": model["name"],
            "provider": model["provider"],
            "tier": model["tier"],
            "task_family": task.primary_family,
            "domain": task.domain,
            "risk_tier": task.risk_tier,
            "workflow_profile": task.workflow_profile,
            "stages": task.required_stages,
            "verifiers_attached": [v["name"] for v in verifier_models],
            "policy_notes": policy.notes,
            "quality_breakdown": {
                "family_fit": model["q"]["family_fit"],
                "workflow_fit": model["q"]["workflow_fit"],
                "domain_fit": model["q"]["domain_fit"],
                "runtime_fit": model["q"]["runtime_fit"],
                "contextual_mean": model["q"]["contextual_mean"],
                "runtime_adjusted_mean": model["q"]["runtime_adjusted_mean"],
            },
            "sla_check": {
                "estimated_latency_ms": est_latency,
                "estimated_cost_usd": est_cost,
                "max_latency_budget_ms": task.request_constraints.max_latency_ms,
                "max_cost_budget_usd": task.request_constraints.max_cost_usd,
                "within_latency_sla": (
                    task.request_constraints.max_latency_ms is None
                    or est_latency <= task.request_constraints.max_latency_ms
                ),
                "within_cost_sla": (
                    task.request_constraints.max_cost_usd is None
                    or est_cost <= task.request_constraints.max_cost_usd
                ),
            }
        },
        trace={
            "router_version": ROUTER_VERSION,
            "parser_version": PARSER_VERSION,
            "policy_version": POLICY_VERSION,
            "scoring_version": SCORING_VERSION,
            "registry_version": REGISTRY_VERSION,
            "timestamp_utc": datetime.now(timezone.utc).isoformat(),
            "confidence_simulations": confidence_data["n_simulations"],
            "win_probabilities": confidence_data["win_probabilities"],
        }
    )


def route(
    prompt: str,
    input_formats: Optional[List[str]] = None,
    estimated_tokens: int = 2000,
    estimated_output_tokens: Optional[int] = None,
    artifact_hints: Optional[List[Dict[str, Any]]] = None,
    llm_parser: Optional[Callable] = None,
    request_constraints: Optional[RequestConstraints] = None,
    tenant_context: Optional[TenantContext] = None,
    files: Optional[List[str]] = None,
    profile_name: str = "balanced",
    registry: Optional[List[Dict[str, Any]]] = None,
    shadow_model: Optional[str] = None,
) -> RoutingDecision:
    """Main entry point for Atlas Router. Returns a complete RoutingDecision."""
    start_time = time.time()
    decision_id = str(uuid.uuid4())
    registry = registry or MODEL_REGISTRY
    # ---- Parse task ----
    task = parse_task_request(
        prompt=prompt,
        input_formats=input_formats,
        estimated_tokens=estimated_tokens,
        estimated_output_tokens=estimated_output_tokens,
        artifact_hints=artifact_hints,
        llm_parser=llm_parser,
        request_constraints=request_constraints,
        tenant_context=tenant_context,
        files = files,
    )
    # ---- Feasibility filtering ----
    feasible, feasibility_reasons = feasibility_filter(task, registry)
    if not feasible:
        return RoutingDecision(
            selected_plan=None, abstain=True, escalate_to_human=True,
            decision_record={
                "decision_id": decision_id,
                "status": "no_feasible_models",
                "feasibility_reasons": feasibility_reasons,
                "task_summary": {
                    "primary_family": task.primary_family,
                    "domain": task.domain,
                    "required_formats": task.required_formats,
                },
                "router_version": ROUTER_VERSION,
                "timestamp_utc": datetime.now(timezone.utc).isoformat(),
                "elapsed_ms": (time.time() - start_time) * 1000,
            }
        )
    # ---- Policy ----
    policy = evaluate_policy(task)
    if policy.must_abstain:
        return RoutingDecision(
            selected_plan=None, abstain=True, escalate_to_human=False,
            decision_record={
                "decision_id": decision_id,
                "status": "policy_abstain",
                "deny_reason": policy.deny_reason,
                "policy_notes": policy.notes,
                "router_version": ROUTER_VERSION,
                "timestamp_utc": datetime.now(timezone.utc).isoformat(),
                "elapsed_ms": (time.time() - start_time) * 1000,
            }
        )
    gated, policy_reasons = apply_policy_to_models(task, policy, feasible)
    if not gated:
        return RoutingDecision(
            selected_plan=None, abstain=True, escalate_to_human=True,
            decision_record={
                "decision_id": decision_id,
                "status": "no_models_after_policy",
                "feasibility_reasons": feasibility_reasons,
                "policy_reasons": policy_reasons,
                "policy_notes": policy.notes,
                "router_version": ROUTER_VERSION,
                "timestamp_utc": datetime.now(timezone.utc).isoformat(),
                "elapsed_ms": (time.time() - start_time) * 1000,
            }
        )
    # ---- Contextual quality ----
    enriched = attach_contextual_quality(gated, task)
    # ---- Pareto reduction (with top-3 floor) ----
    frontier = pareto_frontier(enriched, eps=0.02)
    if len(frontier) < 3 and len(enriched) >= 3:
        ranked = sorted(enriched, key=lambda m: m["q"]["runtime_adjusted_mean"], reverse=True)
        frontier = ranked[:3]
    # ---- Utility scoring ----
    effective_profile = choose_effective_profile(task, profile_name)
    scored = compute_utilities(frontier, task, effective_profile)
    # ---- Confidence estimation ----
    confidence_data = estimate_confidence(scored, task, effective_profile)
    # ---- Select primary and fallbacks ----
    primary_model = scored[0]
    fallback_models = scored[1:4] if len(scored) > 1 else []
    # ---- Verifier planning ----
    verifier_models = choose_verifier_models(
        all_models=feasible,
        verifier_types=policy.require_verifier_types,
        selected_model_name=primary_model["name"],
        max_verifiers=2
    )
    # ---- Plan generation: single-model, then consider multi-stage ----
    single_plan = generate_single_model_plan(
        model=primary_model, task=task, confidence_data=confidence_data,
        verifier_models=verifier_models, policy=policy,
        fallback_models=fallback_models, profile_name=effective_profile,
    )
    plan = single_plan
    if not task.request_constraints.must_use_single_model and len(task.required_stages) >= 2:
        multi_plan = generate_multi_stage_plan(
            models=scored, task=task, confidence_data=confidence_data,
            verifier_models=verifier_models, policy=policy,
            profile_name=effective_profile,
        )
        # Adopt multi-stage only if it clearly improves expected quality
        if multi_plan and multi_plan.expected_quality >= single_plan.expected_quality + 0.03:
            plan = multi_plan
    # ---- Abstention and escalation ladder ----
    abstain = False
    escalate = policy.must_escalate
    if confidence_data["top_confidence"] < CONFIDENCE_ABSTAIN_THRESHOLD:
        abstain = True
        escalate = True
    elif confidence_data["top_confidence"] < CONFIDENCE_ESCALATE_THRESHOLD:
        escalate = True
    elif confidence_data["top_confidence"] < task.request_constraints.min_confidence:
        escalate = True
    if task.risk_tier == "high" and confidence_data["top_confidence"] < CONFIDENCE_HIGH_THRESHOLD:
        escalate = True
    # ---- Shadow routing hook ----
    shadow_info = None
    if shadow_model:
        shadow_candidates = [m for m in scored if m["name"] == shadow_model]
        if shadow_candidates:
            shadow_info = {
                "shadow_model": shadow_model,
                "shadow_utility": shadow_candidates[0]["u"]["expected_utility"],
                "shadow_quality": shadow_candidates[0]["q"]["runtime_adjusted_mean"],
            }
    # ---- Decision record ----
    elapsed_ms = (time.time() - start_time) * 1000
    decision_record = {
        "decision_id": decision_id,
        "status": "routed" if not abstain else "abstained",
        "router_version": ROUTER_VERSION,
        "timestamp_utc": datetime.now(timezone.utc).isoformat(),
        "elapsed_ms": elapsed_ms,
        "task_summary": {
            "primary_family": task.primary_family,
            "secondary_families": task.secondary_families,
            "domain": task.domain,
            "risk_tier": task.risk_tier,
            "risk_type": task.risk_type,
            "complexity": task.complexity,
            "workflow_profile": task.workflow_profile,
            "required_stages": task.required_stages,
            "expected_output": task.expected_output,
            "requires_json": task.requires_json,
            "requires_ocr": task.requires_ocr,
            "requires_web_search": task.requires_web_search,
            "requires_citations": task.requires_citations,
            "requires_verifier": task.requires_verifier,
            "parser_confidence": task.parser_confidence,
            "conflict_flags": task.conflict_flags,
            "detected_languages": task.detected_languages,
            "total_file_size_bytes": task.total_file_size_bytes,
            "inferred_topics": [a.inferred_topic for a in task.artifacts if a.inferred_topic],
            "ambiguity_score": task.ambiguity_score,
            "decomposition_needed": task.decomposition_needed,
        },
        "pipeline_trace": {
            "registry_models": len(registry),
            "feasible_after_filter": len(feasible),
            "feasibility_exclusions": feasibility_reasons,
            "policy_exclusions": policy_reasons,
            "after_policy": len(gated),
            "after_pareto": len(frontier),
            "scored_candidates": len(scored),
            "profile_used": effective_profile,
            "policy_notes": policy.notes,
            "policy_verifier_types": policy.require_verifier_types,
        },
        "confidence": {
            "top_confidence": confidence_data["top_confidence"],
            "margin": confidence_data["margin"],
            "win_probabilities": confidence_data["win_probabilities"],
            "min_confidence_threshold": task.request_constraints.min_confidence,
        },
        "selected_plan_id": plan.plan_id if plan else None,
        "selected_plan_type": plan.plan_type if plan else None,
        "escalate_to_human": escalate,
        "abstain": abstain,
        "shadow_routing": shadow_info,
        "observability": {
            "routing_latency_ms": elapsed_ms,
            "models_evaluated": len(scored),
            "confidence_sims": confidence_data["n_simulations"],
        },
        "learning_hook": {
            "outcome_pending": True,
            "primary_model": primary_model["name"],
            "task_family": task.primary_family,
            "domain": task.domain,
            "confidence": confidence_data["top_confidence"],
            "note": "Record outcome to update priors via Bayesian learning."
        }
    }
    # ---- Reproducibility hash ----
    hash_input = json.dumps({
        "prompt_hash": hashlib.sha256(prompt.encode()).hexdigest()[:16],
        "task_family": task.primary_family,
        "domain": task.domain,
        "profile": effective_profile,
        "feasible_models": sorted([m["name"] for m in feasible]),
        "registry_version": REGISTRY_VERSION,
    }, sort_keys=True)
    decision_record["reproducibility_hash"] = hashlib.sha256(hash_input.encode()).hexdigest()[:24]
    return RoutingDecision(
        selected_plan=plan,
        abstain=abstain,
        escalate_to_human=escalate,
        decision_record=decision_record,
    )


def record_outcome(
    registry: List[Dict[str, Any]],
    model_name: str,
    task_family: str,
    success: bool,
    quality_score: float = 0.0,
    latency_ms: float = 0.0,
    cost_usd: float = 0.0,
    user_accepted: bool = True,
    safety_flagged: bool = False,
) -> None:
    """Update model evaluation state and Bayesian priors from an observed outcome."""
    for model in registry:
        if model["name"] == model_name:
            ev = model["evaluation"]
            ev["samples"] += 1
            ev["quality_sum"] += quality_score
            ev["latency_sum_ms"] += latency_ms
            ev["cost_sum"] += cost_usd
            ev["last_updated"] = datetime.now(timezone.utc).isoformat()
            if success:
                ev["wins"] += 1
            else:
                ev["losses"] += 1
            if user_accepted:
                ev["user_acceptance_sum"] += 1.0
            if safety_flagged:
                ev["safety_flags"] += 1
            fam_prior = model["priors"]["task_family"].get(task_family)
            if fam_prior:
                if success:
                    fam_prior["alpha"] += 1
                else:
                    fam_prior["beta"] += 1
            break


def format_decision_summary(decision: RoutingDecision) -> str:
    lines = []
    rec = decision.decision_record
    lines.append("=== Atlas Router Decision ===")
    lines.append(f"Decision ID: {rec['decision_id']}")
    lines.append(f"Status: {rec['status'].upper()}")
    lines.append(f"Timestamp: {rec['timestamp_utc']}")
    lines.append(f"Routing latency: {rec['elapsed_ms']:.1f} ms")
    lines.append("")
    ts = rec.get("task_summary", {})
    lines.append("-- Task Analysis --")
    lines.append(f"  Primary family: {ts.get('primary_family', 'N/A')}")
    lines.append(f"  Domain: {ts.get('domain', 'N/A')}")
    lines.append(f"  Risk tier: {ts.get('risk_tier', 'N/A')}")
    lines.append(f"  Complexity: {ts.get('complexity', 'N/A')}")
    lines.append(f"  Workflow profile: {ts.get('workflow_profile', 'N/A')}")
    lines.append(f"  Stages: {ts.get('required_stages', [])}")
    lines.append(f"  Parser confidence: {ts.get('parser_confidence', 0):.2f}")
    if ts.get("conflict_flags"):
        lines.append(f" ! Conflicts: {ts['conflict_flags']}")
    if ts.get("inferred_topics"):
        lines.append(f" Inferred topics: {ts['inferred_topics']}")
    if ts.get("detected_languages"):
        lines.append(f" Languages: {ts['detected_languages']}")
    lines.append("")
    if decision.abstain:
        lines.append("!! DECISION: ABSTAIN")
        if rec.get("deny_reason"):
            lines.append(f"   Deny reason: {rec['deny_reason']}")
        if decision.escalate_to_human:
            lines.append("!! ESCALATION: Human review required")
        lines.append("")
    elif decision.selected_plan:
        plan = decision.selected_plan
        lines.append("-- Selected Plan --")
        lines.append(f"  Plan ID: {plan.plan_id}")
        lines.append(f"  Plan type: {plan.plan_type}")
        if plan.plan_type == "single_model":
            lines.append(f"  Primary model: {plan.selected_model}")
        else:
            for sr in plan.stage_routes:
                lines.append(f"    Stage {sr.stage_id} [{sr.stage_name}]: {sr.selected_model}")
        lines.append(f"  Fallbacks: {plan.fallback_models}")
        lines.append(f"  Verifiers: {plan.verifier_models}")
        lines.append(f"  Profile used: {plan.profile_used}")
        lines.append(f"  Expected quality: {plan.expected_quality:.3f}")
        lines.append(f"  Expected latency: {plan.expected_latency_ms:.0f} ms")
        lines.append(f"  Expected cost: ${plan.expected_cost_usd:.5f}")
        lines.append(f"  Confidence: {plan.confidence:.3f}")
        lines.append(f"  Confidence margin: {plan.confidence_margin:.3f}")
        lines.append("")
        if decision.escalate_to_human:
            lines.append("!! ESCALATION: Human review recommended")
            lines.append("")
        qb = plan.explanation.get("quality_breakdown", {})
        if qb:
            lines.append("-- Quality Breakdown --")
            for k, v in qb.items():
                lines.append(f"  {k}: {v:.3f}")
            lines.append("")
        sla = plan.explanation.get("sla_check", {})
        if sla:
            lines.append("-- SLA Check --")
            lines.append(f"  Within latency SLA: {'PASS' if sla.get('within_latency_sla') else 'FAIL'}")
            lines.append(f"  Within cost SLA: {'PASS' if sla.get('within_cost_sla') else 'FAIL'}")
            lines.append("")
    pt = rec.get("pipeline_trace", {})
    if pt:
        lines.append("-- Pipeline Trace --")
        lines.append(f"  Registry models: {pt.get('registry_models', 0)}")
        lines.append(f"  After feasibility: {pt.get('feasible_after_filter', 0)}")
        lines.append(f"  After policy: {pt.get('after_policy', 0)}")
        lines.append(f"  After Pareto: {pt.get('after_pareto', 0)}")
        lines.append(f"  Policy notes: {pt.get('policy_notes', [])}")
        lines.append("")
    conf = rec.get("confidence", {})
    if conf:
        lines.append("-- Confidence --")
        lines.append(f"  Top confidence: {conf.get('top_confidence', 0):.3f}")
        lines.append(f"  Margin: {conf.get('margin', 0):.3f}")
        wp = conf.get("win_probabilities", {})
        for name, prob in sorted(wp.items(), key=lambda x: x[1], reverse=True):
            lines.append(f"    {name}: {prob:.3f}")
        lines.append("")
    lines.append(f"Reproducibility hash: {rec.get('reproducibility_hash', 'N/A')}")
    lines.append("=================================")
    return "\n".join(lines)


# ============================================================================
# END-TO-END DEMONSTRATION
# ============================================================================
def run_demo():
    print("=" * 80)
    print("ATLAS ROUTER — COMPLETE END-TO-END DEMONSTRATION")
    print("=" * 80)
    print()

        # ---- Scenario 1: Coding Task ----
    print("-" * 80)
    print("SCENARIO 1: Complex coding task")
    print("-" * 80)
    decision_1 = route(
        prompt="Implement a concurrent web crawler in Python with rate limiting, retry logic, and structured JSON output of all crawled pages.",
        input_formats=["text"],
        estimated_tokens=3000,
        estimated_output_tokens=4000,
    )
    print(format_decision_summary(decision_1))
    print()
    
    # ---- Scenario 2: Legal Contract Review ----
    print("-" * 80)
    print("SCENARIO 2: Legal contract review (high-risk, verification required)")
    print("-" * 80)
    decision_2 = route(
        prompt="Review this 45-page NDA contract, identify all obligations and liability clauses, flag non-standard terms, and verify compliance with GDPR. Return structured JSON.",
        input_formats=["pdf", "text"],
        estimated_tokens=85000,
        estimated_output_tokens=8000,
        artifact_hints=[{"format": "pdf", "page_count": 45, "text_density": 0.90, "scan_likelihood": 0.05}],
    )
    print(format_decision_summary(decision_2))
    print()
    
    # ---- Scenario 3: Audio Summarization ----
    print("-" * 80)
    print("SCENARIO 3: Audio meeting summarization")
    print("-" * 80)
    decision_3 = route(
        prompt="Summarize this 45-minute engineering meeting recording. Extract action items and decisions made.",
        input_formats=["audio", "text"],
        estimated_tokens=25000,
        estimated_output_tokens=2000,
        artifact_hints=[{"format": "audio", "audio_duration_sec": 2700, "audio_quality": 0.80}],
    )
    print(format_decision_summary(decision_3))
    print()
    
    # ---- Scenario 4: Policy Denial (Offensive Security) ----
    print("-" * 80)
    print("SCENARIO 4: Offensive security request (policy denial)")
    print("-" * 80)
    decision_4 = route(
        prompt="Write an exploit payload for CVE-2024-1234 that achieves remote code execution on unpatched systems.",
        input_formats=["text"],
        estimated_tokens=1500,
        estimated_output_tokens=3000,
    )
    print(format_decision_summary(decision_4))
    print()
    
    # ---- Scenario 5: Budget-Constrained Customer Support ----
    print("-" * 80)
    print("SCENARIO 5: Budget-constrained customer support summarization")
    print("-" * 80)
    decision_5 = route(
        prompt="Summarize these customer support tickets and identify the top 5 recurring issues. Suggest response templates.",
        input_formats=["text"],
        estimated_tokens=8000,
        estimated_output_tokens=3000,
        request_constraints=RequestConstraints(max_latency_ms=3000, max_cost_usd=0.05),
        tenant_context=TenantContext(
            tenant_id="acme-corp", tenant_name="Acme Corporation", budget_remaining_usd=2.50,
        ),
    )
    print(format_decision_summary(decision_5))
    print()
    
    # ---- Scenario 6: Long-context research with citations ----
    print("-" * 80)
    print("SCENARIO 6: Long-context research synthesis with citations required")
    print("-" * 80)
    decision_6 = route(
        prompt="Analyze these research papers on transformer architectures, synthesize the key findings, compare methodologies, and provide citations for each claim.",
        input_formats=["pdf", "text"],
        estimated_tokens=350000,
        estimated_output_tokens=12000,
        artifact_hints=[{"format": "pdf", "page_count": 120, "text_density": 0.92}],
    )
    print(format_decision_summary(decision_6))
    print()
    
    # ---- Scenario 7: File-driven routing with conflict detection ----
    print("-" * 80)
    print("SCENARIO 7: File-driven input (auto-detect + conflict + topic inference)")
    print("-" * 80)
    
    # Create a small temp PDF-like text file to exercise the file path.
    import tempfile
    tmp = tempfile.NamedTemporaryFile(suffix=".txt", delete=False, mode="w")
    tmp.write("THIS AGREEMENT is made between the parties. "
              "The obligations and liability of each party of the first part...")
    tmp.close()
    
    decision_7 = route(
        prompt="Explain this image in detail",   # deliberate conflict: says image, file is text/contract
        files=[tmp.name],
        estimated_tokens=5000,
        estimated_output_tokens=2000,
    )
    print(format_decision_summary(decision_7))
    os.unlink(tmp.name)
    
    
    # ---- Summary table ----
    print("-" * 80)
    print("ROUTING SUMMARY TABLE")
    print("-" * 80)
    summary_data = []
    for i, (desc, dec) in enumerate([
        ("Coding task", decision_1),
        ("Legal contract review", decision_2),
        ("Audio summarization", decision_3),
        ("Offensive security", decision_4),
        ("Budget customer support", decision_5),
        ("Long-context research", decision_6),
        ("File-driven input", decision_7),
    ], 1):
        plan = dec.selected_plan
        if plan and plan.plan_type == "multi_stage":
            model_label = " + ".join(sorted({sr.selected_model for sr in plan.stage_routes}))
        elif plan:
            model_label = plan.selected_model
        else:
            model_label = "NONE"
        summary_data.append({
            "Scenario": f"{i}. {desc}",
            "Selected Model": model_label,
            "Confidence": f"{plan.confidence:.3f}" if plan else "N/A",
            "Est. Cost": f"${plan.expected_cost_usd:.5f}" if plan else "N/A",
            "Est. Latency": f"{plan.expected_latency_ms:.0f}ms" if plan else "N/A",
            "Abstain": dec.abstain,
            "Escalate": dec.escalate_to_human,
            "Verifiers": len(plan.verifier_models) if plan else 0,
        })
    # Simple custom string formatting table instead of pandas
    headers = ['Scenario', 'Selected Model', 'Conf', 'Est. Cost', 'Est. Latency', 'Abstain', 'Escalate', 'Verifiers']
    col_widths = [28, 40, 6, 11, 14, 7, 8, 9]
    
    header_str = ' | '.join(h.ljust(w) for h, w in zip(headers, col_widths))
    print(header_str)
    print('-' * len(header_str))
    
    for row in summary_data:
        row_vals = [
            str(row['Scenario']), str(row['Selected Model']), str(row['Confidence']),
            str(row['Est. Cost']), str(row['Est. Latency']), str(row['Abstain']),
            str(row['Escalate']), str(row['Verifiers'])
        ]
        row_str = ' | '.join(v.ljust(w) for v, w in zip(row_vals, col_widths))
        print(row_str)
    print()
    print("=" * 80)
    print("END OF DEMONSTRATION")
    print("=" * 80)
    
    
    
    

def process_external_feedback(model_name: str, task_family: str, win: bool):
    """External interface for routing feedback to update Bayesian priors."""
    eval_dict = {"wins": 1 if win else 0, "losses": 0 if win else 1, "samples": 1}
    for m in ATLAS_MODEL_REGISTRY:
        if m["name"] == model_name:
            update_model_evaluation(m, task_family, eval_dict)
            break
if __name__ == '__main__':
    run_demo()